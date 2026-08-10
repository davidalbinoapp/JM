"""
Motor de geração do Jornal Mural (AGA).

Respeita a estrutura fixa do PPT:
  1. CAPA           -> slide 0, editado IN PLACE (nunca duplicado/removido)
  2. MATÉRIAS       -> bloco dinâmico, gerado a partir do briefing (dupla/única)
  3. CARTAZES       -> slides fixos, NUNCA tocados
  4. PÁGINA FINAL   -> slide fixo, NUNCA tocado nesta versão (thumbnail
                       reconstruído fica pra uma próxima etapa)
"""
import copy
import io
import os
import re
import math
import functools as _functools
import unicodedata
import zipfile
import tempfile

# Silencia o ruído do OpenCV no terminal (ex.: o WARN "setPreferableTarget
# Targets are not supported by the new graph engine" do OpenCV 5 ao rodar o
# YuNet). Precisa ser ANTES do 1º "import cv2" — como o cv2 só é importado
# dentro das funções (mais pra baixo), setar aqui no topo do módulo garante
# isso. É só aviso, não erro: o enquadramento funciona igual. `setdefault`
# respeita quem já tiver definido a variável por fora.
os.environ.setdefault("OPENCV_LOG_LEVEL", "ERROR")

from PIL import Image
from pptx import Presentation
from pptx.util import Emu


def extract_media_from_template(template_path):
    """Extrai a pasta ppt/media do próprio template pra uma pasta temporária,
    assim o motor não depende de uma pasta de mídia preparada à parte —
    o .pptx do template já carrega a biblioteca de tags de editoria."""
    tmp_dir = tempfile.mkdtemp(prefix="jm_media_")
    with zipfile.ZipFile(template_path) as z:
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                z.extract(name, tmp_dir)
    return os.path.join(tmp_dir, "ppt", "media")

import os as _os

TAGS_DIR = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "tags_editorias")

EDITORIA_TAGS = {
    "SEGURANÇA": "seguranca.png",
    "NOSSA EMPRESA": "nossa_empresa.png",
    "LEGADO": "legado.png",
    "EVENTOS": "eventos.png",
    "MEIO AMBIENTE": "meio_ambiente.png",
    "DIVERSIDADE": "diversidade.png",
    "ESSÊNCIA AGA": "essencia_aga.png",
    "SAÚDE E BEM-ESTAR": "saude_e_bem_estar.png",
    "INOVAÇÃO E TECNOLOGIA": "inovacao_e_tecnologia.png",
    "PESSOAS": "pessoas.png",
    "EXCELÊNCIA OPERACIONAL": "excelencia_operacional.png",
}

FONTS_DIR = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "fonts")
MODELS_DIR = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "models")
YUNET_PATH = _os.path.join(MODELS_DIR, "face_detection_yunet.onnx")

# a fonte real do template é Arial. Ela não pode ser redistribuída aqui
# (é uma fonte comercial da Monotype), mas Mac e Windows normalmente já
# vêm com ela instalada — então tentamos usar a Arial de verdade primeiro,
# e só caímos pra Liberation Sans (compatível em métrica com a Arial,
# feita de propósito pra isso) se a Arial não for encontrada no sistema.
_CANDIDATOS_ARIAL_REGULAR = [
    "/Library/Fonts/Arial.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "C:\\Windows\\Fonts\\arial.ttf",
    "/usr/share/fonts/truetype/msttcorefonts/Arial.ttf",
]
_CANDIDATOS_ARIAL_BOLD = [
    "/Library/Fonts/Arial Bold.ttf",
    "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
    "C:\\Windows\\Fonts\\arialbd.ttf",
    "/usr/share/fonts/truetype/msttcorefonts/Arial_Bold.ttf",
]


def _fonte_abre(caminho):
    """Só aceita a fonte se o PIL conseguir abrir de fato. Existir o arquivo
    não basta: pode ser um .ttf quebrado, um atalho, ou um arquivo baixado
    pela metade — e aí o erro só apareceria lá na frente, no meio da
    montagem, como um 'cannot open resource' sem explicação."""
    from PIL import ImageFont
    try:
        ImageFont.truetype(caminho, 40)
        return True
    except Exception:
        return False


def _resolver_caminho_fonte(negrito):
    reserva = _os.path.join(FONTS_DIR, "LiberationSans-Bold.ttf" if negrito
                            else "LiberationSans-Regular.ttf")
    for caminho in (_CANDIDATOS_ARIAL_BOLD if negrito else _CANDIDATOS_ARIAL_REGULAR):
        if _os.path.isfile(caminho) and _fonte_abre(caminho):
            return caminho
    if _os.path.isfile(reserva) and _fonte_abre(reserva):
        return reserva
    raise RuntimeError(
        "Nenhuma fonte utilizável foi encontrada. O sistema procura a Arial "
        "instalada no computador e, se não achar, usa a Liberation Sans que "
        "vem junto do projeto. Confira se a pasta 'fonts/' está do lado do "
        "jm_engine.py com os arquivos LiberationSans-Regular.ttf e "
        "LiberationSans-Bold.ttf (cerca de 400 KB cada)."
    )


FONT_REGULAR = _resolver_caminho_fonte(negrito=False)
FONT_BOLD = _resolver_caminho_fonte(negrito=True)


@_functools.lru_cache(maxsize=80000)
def _contar_linhas_texto(texto, largura_emu, tamanho_pt, negrito=False):
    """Estima quantas linhas o texto vai ocupar dentro de uma caixa de
    largura `largura_emu`, pra permitir reposicionar elementos abaixo dela
    de forma proporcional ao conteúdo real (em vez de uma posição fixa).
    Cacheado e medindo via _medir_texto (o balanceador conta milhares de wraps)."""
    import math
    px_size = max(int(tamanho_pt * 4), 8)
    # desconta as margens internas padrão da caixa de texto do PowerPoint
    # (0,1" de cada lado). Sem isso o cálculo usa a largura cheia e acha que
    # cabe mais texto por linha do que realmente cabe — contando uma linha a
    # menos em títulos que ficam no limiar, e deixando o corpo subir demais.
    MARGEM_INTERNA_EMU = 91440  # 0,1 polegada
    largura_util = max(largura_emu - 2 * MARGEM_INTERNA_EMU, 12700)
    largura_px = (largura_util / 12700) * 4

    def _largura(s):
        return _medir_texto(s, px_size, negrito)

    def _linhas_do_token(s):
        # Um token indivisível (uma palavra longa, ou um par grudado pelo
        # espaço inseparável da prevenção de viúva) que seja MAIS LARGO que a
        # caixa é quebrado pelo PowerPoint mesmo sem ponto de quebra. Aqui a
        # gente conta quantas linhas ele realmente ocupa. Sem isso o contador
        # trata esse token como 1 linha só, subestima a altura do título e o
        # corpo é posicionado alto demais, invadindo o título.
        return max(1, math.ceil(_largura(s) / largura_px))

    total_linhas = 0
    for paragrafo in texto.split("\n"):
        # split só no espaço normal: o espaço inseparável ( ) usado pra
        # evitar viúva NÃO é ponto de quebra. Se usarmos o .split() padrão do
        # Python (que trata   como espaço), o contador acha que cabe uma
        # linha a menos do que o PowerPoint realmente quebra — e aí o corpo é
        # posicionado alto demais e invade o título.
        palavras = [p for p in paragrafo.split(" ") if p]
        if not palavras:
            total_linhas += 1  # linha em branco também ocupa altura
            continue
        linha_atual = palavras[0]
        linhas_neste_paragrafo = _linhas_do_token(palavras[0])
        for palavra in palavras[1:]:
            teste = linha_atual + " " + palavra
            if _largura(teste) > largura_px:
                # quebra: a palavra desce pra próxima linha (e, se ela mesma
                # for mais larga que a caixa, ocupa mais de uma linha)
                linhas_neste_paragrafo += _linhas_do_token(palavra)
                linha_atual = palavra
            else:
                linha_atual = teste
        total_linhas += linhas_neste_paragrafo
    return max(total_linhas, 1)


def _fonte_pt_do_titulo(title_box, default=40):
    """Tamanho (pt) da fonte atual do título, pra medir quebra de linha."""
    runs = title_box.text_frame.paragraphs[0].runs
    if runs and runs[0].font.size:
        return runs[0].font.size / 12700
    return default


def _fonte_pt_do_corpo(body_box, default=19.0):
    """Tamanho (pt) da fonte do corpo — o tamanho de leitura do template."""
    runs = body_box.text_frame.paragraphs[0].runs
    if runs and runs[0].font.size:
        return runs[0].font.size / 12700
    return default


# quanto, no máximo, dá pra reduzir a fonte do título pra tirar uma linha (12%).
# Reduções maiores que isso descaracterizariam o título — aí prefere-se manter o
# tamanho do template e deixar a linha a mais.
FATOR_MIN_FONTE_TITULO = 0.88

# Encaixe título x corpo: o CORPO é prioridade (fica no tamanho de leitura). Quem
# cede primeiro é o TÍTULO — pode encolher até 15% (ex.: 40->34) e aceitar uma
# viúva, que num título com muito texto não tem problema. Só se, mesmo com o
# título no piso e a caixa alargada, o corpo NÃO couber, ele desce de 19 até no
# MÁXIMO 17pt (leitura garantida). O 12pt de antes era um bug e nunca mais pode
# acontecer. Ver _encaixar_corpo_preservando.
FATOR_MIN_TITULO_ENCAIXE = 0.85
PT_PISO_CORPO = 17.0

# Balanceamento da MASSA de corpo (a "percepção" do David: olhar o bloco inteiro e
# estreitar/alargar um tico a caixa até o texto distribuir bem, sem linha curta no
# meio deixando um vazão — ex.: "...e o [VÃO] / desenvolvimento das comunidades").
# Não é largura fixa: é modular por texto, dentro de um limite, sempre prezando o
# resultado. Ver _encaixar_corpo_preservando.
_DEBUG_MASSA = _os.environ.get("JM_DEBUG_MASSA") == "1"
# A nossa medição (PIL) fecha MAIS texto por linha do que o PowerPoint/LibreOffice
# realmente fecham. Sem corrigir isso, o balanceador otimiza um wrap que não é o
# real (escolhe uma largura que "no cálculo" fecha em 4 linhas limpas, mas na tela
# vira 5 com um vão). FATOR_WRAP_RENDER aperta a medição pra bater com a régua real
# — calibrado pro renderizador mais apertado (o PowerPoint do David).
FATOR_WRAP_RENDER = 0.955
NARROW_LIM_MASSA = 0.12       # estreita no máx 12% da largura de projeto
ALVO_ULTIMA_MASSA = 0.35      # última linha abaixo disso = viúva -> penaliza
PESO_ULTIMA_MASSA = 1.5       # peso da penalidade de viúva de fim
MARGEM_MASSA = 0.05           # só troca da config atual se melhorar mais que isso
                              # (viés anti-churn: página já boa não muda)


@_functools.lru_cache(maxsize=200000)
def _medir_texto(s, px_fonte, negrito):
    """Largura (px) de um texto na fonte do JM, cacheada (o balanceador mede
    milhares de wraps)."""
    from PIL import ImageFont
    fonte = FONT_BOLD if negrito else FONT_REGULAR
    f = ImageFont.truetype(fonte, max(int(px_fonte), 8))
    b = f.getbbox(s)
    return b[2] - b[0]


def _fills_corpo(texto, largura_emu, pt, negrito=False, fator=1.0):
    """Fração de preenchimento de cada linha do corpo (quebra gulosa, \xa0 é
    inquebrável). `fator`<1 modela o PowerPoint (mais apertado). Serve pra medir
    a qualidade da distribuição da massa (linhas curtas no meio = vazão)."""
    MARG = 91440
    px = (max(largura_emu - 2 * MARG, 12700) / 12700) * 4 * fator
    pxf = int(pt * 4)
    toks = [t for t in texto.split(" ") if t]
    if not toks:
        return []
    linhas = [toks[0]]
    for t in toks[1:]:
        if _medir_texto(linhas[-1] + " " + t, pxf, negrito) > px:
            linhas.append(t)
        else:
            linhas[-1] = linhas[-1] + " " + t
    return [_medir_texto(l.replace("\xa0", " "), pxf, negrito) / px for l in linhas]


def _reduzir_fonte_titulo_menos_linhas(texto, largura_emu, pt_atual, negrito=True):
    """Se, no tamanho do template, o título quebra em N linhas mas uma redução
    PEQUENA de fonte o encaixa em N-1 (aproveitando o espaço que sobra), devolve
    o MAIOR tamanho inteiro que tira uma linha — nunca abaixo de ~88% do
    original (FATOR_MIN_FONTE_TITULO). Título que já quebra bem, ou que só
    reduziria linha com fonte bem menor, fica com o tamanho original. Só ENCOLHE,
    nunca cresce — então nada de regressão em quem já estava bom."""
    if not texto or not largura_emu:
        return pt_atual
    pt_min = max(34.0, round(pt_atual * FATOR_MIN_FONTE_TITULO))
    linhas0 = _contar_linhas_texto(texto, largura_emu, pt_atual, negrito)
    if linhas0 <= 1:
        return pt_atual
    pt = pt_atual
    while pt - 1 >= pt_min:
        pt -= 1
        if _contar_linhas_texto(texto, largura_emu, pt, negrito) < linhas0:
            return float(pt)   # descendo de 1 em 1, o 1º que tira linha é o MAIOR
    return pt_atual


def _evitar_viuva(texto, largura_emu=None, tamanho_pt=None, negrito=True):
    """Troca o último espaço por espaço inseparável, pra a última palavra
    nunca ficar sozinha numa linha (viúva tipográfica). Vale pro título e pra
    massa de texto do corpo (`negrito=False` mede na fonte regular do corpo).

    Se a largura da caixa e o tamanho da fonte forem informados, só gruda
    as duas últimas palavras quando elas CABEM juntas numa linha. Grudar um
    par que não cabe é contraproducente: o PowerPoint quebra o bloco
    inseparável mesmo assim, sobrando um caco (ex.: um "?" sozinho) — pior
    que a viúva que a gente queria evitar. Nesse caso, deixa quebrar natural."""
    palavras = texto.split(" ")
    if len(palavras) < 2:
        return texto
    if largura_emu and tamanho_pt:
        from PIL import ImageFont
        MARGEM_INTERNA_EMU = 91440  # 0,1" de cada lado, como no contador de linhas
        largura_util = max(largura_emu - 2 * MARGEM_INTERNA_EMU, 12700)
        largura_px = (largura_util / 12700) * 4
        fonte_path = FONT_BOLD if negrito else FONT_REGULAR
        font = ImageFont.truetype(fonte_path, max(int(tamanho_pt * 4), 8))
        par = palavras[-2] + " " + palavras[-1]
        b = font.getbbox(par)
        if (b[2] - b[0]) > largura_px:
            return texto  # nao cabe grudado: deixa quebrar natural
    return " ".join(palavras[:-1]) + " " + palavras[-1]

def _evitar_viuva_corpo(body_box, texto):
    """Equilibra as DUAS ÚLTIMAS linhas do corpo. Em vez de só evitar a viúva de
    uma palavra (descer uma), desce quantas palavras forem necessárias pra a
    última linha ficar parelha com a penúltima ("faço uma média"). As palavras
    que descem são grudadas com espaço inseparável. Regras de segurança: nunca
    aumenta o nº de linhas; nunca deixa a última MAIOR que a penúltima — senão
    ENCALHARIA a penúltima numa linha curta (ex.: "baixo custo." sozinho, com o
    bloco grande embaixo); varre TODOS os k e pega o de menor diferença de
    largura (não para no 1º platô, que não é monotônico). A medição usa a
    largura cheia da caixa (sem "aperto"): o corte tem que casar com o wrap real
    do arquivo, senão a decisão fica boa no cálculo e ruim na página."""
    runs = body_box.text_frame.paragraphs[0].runs
    pt = (runs[0].font.size / 12700) if (runs and runs[0].font.size) else 19
    return _glue_viuva_corpo(texto, body_box.width, pt)


@_functools.lru_cache(maxsize=40000)
def _glue_viuva_corpo(texto, largura_emu, pt, fator=1.0):
    """Núcleo PURO da prevenção de viúva do corpo: devolve o texto com as N
    últimas palavras grudadas por espaço inseparável. Usado pra SETAR o corpo E
    pra MEDIR a distribuição real no balanceador (que precisa enxergar o texto
    exato que vai pro slide, com o nbsp, não o cru). `fator`<1 aperta a régua pra
    bater com o renderizador real. Cacheado e medindo via _medir_texto (também
    cacheado) — o balanceador chama isso milhares de vezes."""
    palavras = [w for w in texto.split(" ") if w]
    if len(palavras) < 3:
        return texto
    MARGEM_INTERNA_EMU = 91440
    largura_px = (max(largura_emu - 2 * MARGEM_INTERNA_EMU, 12700) / 12700) * 4 * fator
    pxf = int(pt * 4)

    def _w(s):
        return _medir_texto(s, pxf, False)

    def _sim(k):
        # junta as últimas k palavras num token só e simula a quebra gulosa;
        # retorna (largura_penúltima, largura_última, nº_de_linhas) ou None
        juntas = " ".join(palavras[-k:])
        if _w(juntas) > largura_px:
            return None
        tokens = palavras[:-k] + [juntas]
        linhas = [tokens[0]]
        for t in tokens[1:]:
            if _w(linhas[-1] + " " + t) > largura_px:
                linhas.append(t)
            else:
                linhas[-1] = linhas[-1] + " " + t
        if len(linhas) < 2:
            return None
        return _w(linhas[-2]), _w(linhas[-1]), len(linhas)

    base = _sim(2)                 # k=2 = a prevenção de viúva simples
    if base is None:
        return texto
    n_linhas = base[2]
    melhor_k, melhor_dif = 2, abs(base[0] - base[1])
    for k in range(3, len(palavras)):
        r = _sim(k)
        if r is None:              # o bloco não cabe mais numa linha -> para
            break
        if r[2] != n_linhas:       # mudaria o nº de linhas -> ignora esse k
            continue
        if r[1] > r[0]:            # última ficaria MAIOR que a penúltima ->
            continue               # encalharia a penúltima curta; ignora
        if abs(r[0] - r[1]) < melhor_dif:   # menor desequilíbrio (global)
            melhor_k, melhor_dif = k, abs(r[0] - r[1])
    # junta as que descem com espaço INSEPARÁVEL pra ficarem na mesma última
    # linha; a separação com o resto é espaço normal (ponto de quebra)
    return " ".join(palavras[:-melhor_k]) + " " + " ".join(palavras[-melhor_k:])


# Aperto da caixa de corpo pra equilibrar a última linha (a "decisão" que a
# gente toma no olho: fechar um tico a caixa até a frase de fecho cair inteira).
CAP_APERTO_CORPO = 0.08          # aperta no máx 8% da largura; além disso, desiste
LIMIAR_WIDOW_CORPO = 0.55        # última linha < disso no wrap natural = "widow" -> tenta apertar
ALVO_ULTIMA_CORPO = 0.65         # depois do aperto, última linha tem que encher pelo menos isso


def _apertar_caixa_corpo(body_box, texto):
    """Aperta LEVEMENTE a largura da caixa de corpo pra tirar uma 'viúva' de fim
    de parágrafo. Quando a última linha fica curta (ex.: 'equipes premiadas!'
    sozinho, com um 'às' pendurado na linha de cima), estreitar um pouco a caixa
    faz o parágrafo refluir e a frase de fecho cair inteira numa linha só —
    exatamente o que a gente faz no olho, fechando um tico a caixa até equilibrar.

    Análise (não chute): varre a largura pra baixo e escolhe a MENOR redução
    (≤ CAP_APERTO_CORPO) que enche a última linha (≥ ALVO_ULTIMA_CORPO) SEM
    aumentar o nº de linhas. Só age se existir mesmo uma última linha curta
    (< LIMIAR_WIDOW_CORPO) — corpo que já fecha bem fica INTACTO (zero regressão).
    Não achando nada dentro do teto, a caixa não muda. Mexe só na largura; como o
    nº de linhas é preservado, o topo/altura do corpo não mudam e não há
    risco de transbordo. Mede na largura cheia da caixa (o mesmo wrap real do
    arquivo), como o resto do motor."""
    from PIL import ImageFont
    runs = body_box.text_frame.paragraphs[0].runs
    pt = (runs[0].font.size / 12700) if (runs and runs[0].font.size) else 19
    palavras = [w for w in texto.split(" ") if w]
    if len(palavras) < 4:
        return
    MARGEM_INTERNA_EMU = 91440
    font = ImageFont.truetype(FONT_REGULAR, max(int(pt * 4), 8))

    def _w(s):
        b = font.getbbox(s)
        return b[2] - b[0]

    def _px(width_emu):
        return (max(width_emu - 2 * MARGEM_INTERNA_EMU, 12700) / 12700) * 4

    def _wrap(largura_px):
        linhas = [palavras[0]]
        for p in palavras[1:]:
            if _w(linhas[-1] + " " + p) > largura_px:
                linhas.append(p)
            else:
                linhas[-1] = linhas[-1] + " " + p
        return linhas

    largura0 = body_box.width
    base = _wrap(_px(largura0))
    n0 = len(base)
    if n0 < 2:
        return
    if (_w(base[-1]) / _px(largura0)) >= LIMIAR_WIDOW_CORPO:
        return  # última linha já está cheia o bastante -> não mexe
    passo = int(0.05 * 360000)            # varre de 0,05cm em 0,05cm
    largura_min = int(largura0 * (1 - CAP_APERTO_CORPO))
    largura = largura0 - passo
    while largura >= largura_min:
        linhas = _wrap(_px(largura))
        if len(linhas) == n0 and (_w(linhas[-1]) / _px(largura)) >= ALVO_ULTIMA_CORPO:
            body_box.width = largura       # menor redução que equilibra
            return
        largura -= passo
    # nada dentro do teto resolveu -> caixa intacta


# O PowerPoint quebra ~3% mais apertado que a nossa medição (PIL/LibreOffice).
# Por isso uma última linha que aqui fica composta pode, no PP, perder a palavra
# da frente e virar um toco solto (a viúva "(NR) nº22", com "Regulamentadora"
# pendurado na linha de cima). FATOR_WRAP_PP modela esse aperto; ALVO_ULTIMA_PP é
# o quanto a última linha precisa encher (no modelo PP) pra NÃO ser um toco.
FATOR_WRAP_PP = 0.97
ALVO_ULTIMA_PP = 0.42


def _compor_ultima_linha_pp(body_box):
    """Estreita LEVEMENTE a caixa do corpo até a última linha parar de ser um
    toco solto QUANDO O POWERPOINT quebra. Espelha o ajuste manual do David na
    capa: em vez de compensar no texto (espaço inseparável), encolhe um tico a
    caixa até a palavra companheira descer junto com o fim. Mede a quebra no
    modelo PP (mais apertado), com o espaço inseparável tratado como token único.
    Bounded por CAP_APERTO_CORPO; se nada dentro do teto compõe, deixa como está.
    Idempotente: se a última linha já está composta no modelo PP, não mexe."""
    from PIL import ImageFont
    runs = body_box.text_frame.paragraphs[0].runs
    if not runs:
        return
    pt = (runs[0].font.size / 12700) if runs[0].font.size else 19
    texto = body_box.text_frame.text          # já com o espaço inseparável, se houver
    MARG_INT = 91440
    font = ImageFont.truetype(FONT_REGULAR, max(int(pt * 4), 8))

    def _w(s):
        b = font.getbbox(s)
        return b[2] - b[0]

    def _fill_ultima(largura_emu, fator):
        # quebra gulosa; fator<1 modela o PP (mais apertado). \xa0 fica no token.
        px = (max(largura_emu - 2 * MARG_INT, 12700) / 12700) * 4 * fator
        toks = [t for t in texto.split(" ") if t]
        if not toks:
            return 1.0
        linhas = [toks[0]]
        for t in toks[1:]:
            if _w(linhas[-1] + " " + t) > px:
                linhas.append(t)
            else:
                linhas[-1] = linhas[-1] + " " + t
        if len(linhas) < 2:
            return 1.0
        return _w(linhas[-1].replace("\xa0", " ")) / px

    def _composta(largura_emu):
        # última linha "composta" (companheira junto, não é toco) nos DOIS
        # renderizadores: no nosso (LibreOffice/PIL, fator 1.0) E no PP (0.97).
        # Exigir os dois evita cair numa largura azarada boa só num deles.
        return min(_fill_ultima(largura_emu, 1.0),
                   _fill_ultima(largura_emu, FATOR_WRAP_PP)) >= ALVO_ULTIMA_PP

    largura0 = body_box.width
    # se já está composta nos dois modelos, não mexe
    if _fill_ultima(largura0, FATOR_WRAP_PP) >= ALVO_ULTIMA_PP:
        return
    passo = int(0.05 * 360000)
    largura_min = int(largura0 * (1 - CAP_APERTO_CORPO))
    largura = largura0 - passo
    while largura >= largura_min:
        if _composta(largura):              # composta nos DOIS modelos -> ok
            body_box.width = largura
            return
        largura -= passo
    # nada dentro do teto compôs nos dois -> caixa intacta (não piora)


# Abaixo desta fração da largura da caixa, uma linha do título que NÃO é a
# última é considerada "órfã" (curta demais no meio do título) e dispara o
# rebalanceamento. Acima disso, a quebra gulosa é tida como boa e mantida.
LIMIAR_ORFAO_TITULO = 0.55
# Fim curto num título de 2 linhas: se a ÚLTIMA linha é um toco muito curto
# (viúva tipo "de FPS", com a 1ª linha quase cheia), reequilibra as duas linhas.
# Bem baixo de propósito: só pega o toco EXTREMO, não a 2ª linha naturalmente
# mais curta que quase todo título tem (senão reestruturaria meia publicação).
LIMIAR_TOCO_TITULO = 0.30


def _quebrar_titulo_balanceado(texto, largura_emu, tamanho_pt, negrito=True):
    """Reequilibra as quebras do título: mantém o MESMO número de linhas que a
    quebra natural (gulosa) do PowerPoint, mas distribui as palavras pra
    minimizar a largura da linha mais cheia — evita que uma palavra curta
    ('CMG') fique órfã sozinha enquanto outra linha fica lotada. Como cada linha
    fica MENOS cheia que na gulosa, nunca transborda mais do que já transbordava
    (a maior linha do balanceado ≤ a maior linha da gulosa, que já cabia).
    Devolve o título com '\\n' nos pontos de quebra; inalterado se cabe em 1
    linha ou faltar largura/fonte. Preserva o espaço inseparável da viúva (o
    par grudado continua um token só)."""
    if not (largura_emu and tamanho_pt):
        return texto
    try:
        from PIL import ImageFont
        MARG = 91440  # 0,1" de cada lado, como no contador de linhas
        largura_px = (max(largura_emu - 2 * MARG, 12700) / 12700) * 4
        fonte = FONT_BOLD if negrito else FONT_REGULAR
        font = ImageFont.truetype(fonte, max(int(tamanho_pt * 4), 8))

        def W(s):
            b = font.getbbox(s)
            return b[2] - b[0]

        tokens = [t for t in texto.split(" ") if t]  # \xa0 (viúva) fica no token
        if len(tokens) < 2:
            return texto
        # quebra gulosa (o que o PowerPoint faria sozinho): guarda as linhas
        greedy = []
        cur = tokens[0]
        for t in tokens[1:]:
            if W(cur + " " + t) > largura_px:
                greedy.append(cur)
                cur = t
            else:
                cur = cur + " " + t
        greedy.append(cur)
        n = len(greedy)
        if n < 2:
            return texto
        # GATE: rebalanceia se a quebra gulosa deixou (a) uma linha curta NO MEIO
        # — o órfão tipo "CMG" —, OU (b) num título de 2 linhas, um TOCO muito
        # curto no FIM (viúva tipo "de FPS", com a 1ª linha quase cheia): aí o
        # balanceado joga uma palavra pra baixo e as duas linhas fecham parelhas.
        # Se o guloso já está bem distribuído, não mexe (nada de regressão).
        meio_curto = any(W(l) < LIMIAR_ORFAO_TITULO * largura_px for l in greedy[:-1])
        fim_toco = (n == 2 and W(greedy[-1]) < LIMIAR_TOCO_TITULO * largura_px)
        if not (meio_curto or fim_toco):
            return texto
        m = len(tokens)
        # dp[k][i] = menor "largura da maior linha" ao quebrar tokens[i:] em k linhas
        INF = float("inf")
        dp = [[INF] * (m + 1) for _ in range(n + 1)]
        nxt = [[m] * (m + 1) for _ in range(n + 1)]
        dp[0][m] = 0.0
        for k in range(1, n + 1):
            for i in range(m - 1, -1, -1):
                # a 1ª linha leva tokens[i:j]; o resto (tokens[j:]) em k-1 linhas.
                # j vai até deixar pelo menos (k-1) tokens pro resto.
                for j in range(i + 1, m - (k - 1) + 1):
                    cand = max(W(" ".join(tokens[i:j])), dp[k - 1][j])
                    if cand < dp[k][i]:
                        dp[k][i] = cand
                        nxt[k][i] = j
        if dp[n][0] == INF:
            return texto
        linhas = []
        i = 0
        for k in range(n, 0, -1):
            j = nxt[k][i]
            linhas.append(" ".join(tokens[i:j]))
            i = j
        return "\n".join(linhas)
    except Exception:
        return texto


def _reposicionar_corpo_apos_titulo(title_box, body_box, texto_titulo=None):
    """Ajusta o topo da caixa de corpo com base na altura REAL que o título
    vai ocupar (em vez de confiar numa posição fixa do template, que só
    funciona se o título tiver exatamente o mesmo número de linhas do
    título original usado no design).

    Conta o texto que ESTÁ na caixa (com o espaço inseparável da viúva já
    aplicado e com a largura/fonte já ajustadas por causa do QR), e não o
    título original — senão a contagem quebra em um número de linhas
    diferente do que o PowerPoint realmente mostra, e o corpo acaba
    posicionado alto demais, invadindo o título."""
    run = title_box.text_frame.paragraphs[0].runs[0]
    tamanho_pt = (run.font.size / 12700) if run.font.size else 28
    texto_real = title_box.text_frame.text or texto_titulo or ""
    linhas = _contar_linhas_texto(texto_real, title_box.width, tamanho_pt, negrito=True)
    altura_linha_emu = int(tamanho_pt * 12700 * 1.25)
    gap_padrao_emu = int(0.5 * 360000)
    novo_top = title_box.top + linhas * altura_linha_emu + gap_padrao_emu
    body_box.top = novo_top


def _ajustar_corpo_sem_transbordar(body_box, texto_corpo, limite_inferior_emu, tamanho_min_pt=11):
    """Depois de reposicionar o corpo, garante que o texto não vaza pro
    espaço do próximo elemento abaixo — reduz o tamanho da fonte
    gradualmente até caber na altura disponível. É isso que impede
    sobreposição quando o título E o corpo são grandes ao mesmo tempo."""
    from pptx.util import Pt
    if not body_box.text_frame.paragraphs[0].runs:
        return
    run = body_box.text_frame.paragraphs[0].runs[0]
    tamanho_pt = (run.font.size / 12700) if run.font.size else 19
    disponivel = limite_inferior_emu - body_box.top - int(0.3 * 360000)
    if disponivel <= 0:
        return

    while tamanho_pt > tamanho_min_pt:
        linhas = _contar_linhas_texto(texto_corpo, body_box.width, tamanho_pt, negrito=False)
        altura_necessaria = linhas * tamanho_pt * 12700 * 1.25
        if altura_necessaria <= disponivel:
            break
        tamanho_pt -= 1

    for paragraph in body_box.text_frame.paragraphs:
        for r in paragraph.runs:
            r.font.size = Pt(tamanho_pt)


def _custo_massa(texto_corpo, largura, pt_b, capa=False):
    """Quão MAL distribuída está a massa de corpo numa dada largura/fonte: soma
    dos 'vazões' (1-preenchimento)² das linhas que NÃO são a última — é isso que
    pega a linha curta no meio com uma palavra longa empurrada pra baixo. Soma
    ainda uma penalidade se a ÚLTIMA linha ficar um toco (viúva). Na capa, mede a
    última linha também no modelo PowerPoint (mais apertado). Menor = melhor.

    IMPORTANTE: mede o texto COM o espaço inseparável já aplicado (o que de fato
    vai pro slide), senão otimizaria um wrap que não é o real — foi o bug do
    "gruda 3 palavras no fim e abre um vão na linha de cima"."""
    # tudo medido na régua REAL (FATOR_WRAP_RENDER): glue e fills no mesmo aperto
    # do renderizador, senão a decisão fica boa no cálculo e ruim na tela.
    texto = _glue_viuva_corpo(texto_corpo, largura, pt_b, FATOR_WRAP_RENDER)
    fills = _fills_corpo(texto, largura, pt_b, fator=FATOR_WRAP_RENDER)
    if not fills:
        return 0.0
    c = sum((1.0 - f) ** 2 for f in fills[:-1])
    if len(fills) >= 2 and fills[-1] < ALVO_ULTIMA_MASSA:
        c += PESO_ULTIMA_MASSA * (ALVO_ULTIMA_MASSA - fills[-1]) ** 2
    return c


def _encaixar_corpo_preservando(title_box, body_box, texto_titulo, texto_corpo,
                                limite_inferior_emu, largura_max_corpo_emu=None,
                                warnings=None, contexto="", capa=False):
    """Encaixa título+corpo espelhando a PERCEPÇÃO do David: olha a massa de texto
    inteira e ESTREITA/ALARGA um tico a caixa (modular por texto, dentro de um
    limite) até distribuir bem — sem linha curta no meio deixando um vazão, nem
    viúva no fim. Ao mesmo tempo mantém o CORPO em tamanho de leitura (19pt); quem
    cede primeiro é o TÍTULO (reduz a fonte, aceita viúva de título); só em último
    caso o corpo desce de 19 até no MÁX 17pt; e se nem a 17 couber, avisa.

    Faz UMA varredura conjunta de (fonte do título) × (largura da caixa) × (fonte
    do corpo 19→18→17), pegando, entre as combinações que CABEM na vertical, a de
    MENOR custo de massa (melhor distribuição), com viés anti-churn: se a config
    atual já está boa, não mexe. Assim o mesmo mecanismo resolve capa, única e
    dupla — vazão, viúva e transbordo tratados como um problema só."""
    from pptx.util import Pt
    if not (title_box and body_box):
        return
    if not (body_box.text_frame.paragraphs[0].runs and title_box.text_frame.paragraphs[0].runs):
        return
    GAP = int(0.5 * 360000)
    FOLGA = int(0.3 * 360000)
    ALTURA_LINHA = 1.25
    PASSO = int(0.05 * 360000)
    pt_tit_atual = _fonte_pt_do_titulo(title_box)
    pt_tit_piso = int(max(30.0, round(pt_tit_atual * FATOR_MIN_TITULO_ENCAIXE)))
    pt_corpo = _fonte_pt_do_corpo(body_box)
    largura0 = body_box.width
    # leque de larguras: estreita até NARROW_LIM_MASSA; alarga até o limite que o
    # layout permite (ex.: embaixo do QR, ou None = só a largura de projeto).
    largura_min = int(largura0 * (1 - NARROW_LIM_MASSA))
    largura_max = int(largura_max_corpo_emu) if largura_max_corpo_emu else largura0
    largura_max = max(largura_max, largura0)
    larguras = list(range(largura_min, largura_max + 1, PASSO))
    if largura0 not in larguras:
        larguras.append(largura0)

    # topo do corpo dado o tamanho do título (só depende de pt_tit): cacheado
    top_por_pt = {}
    def _body_top(pt_tit):
        if pt_tit not in top_por_pt:
            lt = _contar_linhas_texto(texto_titulo, title_box.width, pt_tit, negrito=True)
            top_por_pt[pt_tit] = title_box.top + int(lt * pt_tit * 12700 * ALTURA_LINHA) + GAP
        return top_por_pt[pt_tit]

    plano = None
    transborda = False
    # tenta 19; só cai pra 18 e depois 17 se NADA (título+largura) couber em 19.
    for pt_b in (pt_corpo, 18.0, PT_PISO_CORPO):
        if pt_b > pt_corpo:
            continue
        # PERF: custo e nº de linhas do corpo só dependem de (largura, pt_b), NÃO
        # do título — então mede uma vez por largura, não uma vez por par.
        custo_larg = {}
        altura_larg = {}
        for largura in larguras:
            custo_larg[largura] = _custo_massa(texto_corpo, largura, pt_b, capa=capa)
            n = _contar_linhas_texto(texto_corpo, int(largura * FATOR_WRAP_RENDER), pt_b, negrito=False)
            altura_larg[largura] = int(n * pt_b * 12700 * ALTURA_LINHA)
        candidatos = []
        for pt_tit in range(int(round(pt_tit_atual)), pt_tit_piso - 1, -1):
            disp = limite_inferior_emu - _body_top(pt_tit) - FOLGA
            for largura in larguras:
                if altura_larg[largura] <= disp:      # cabe na vertical
                    chave = (round(custo_larg[largura], 4),
                             int(round(pt_tit_atual)) - pt_tit, abs(largura - largura0))
                    candidatos.append((chave, float(pt_tit), largura))
        if candidatos:
            candidatos.sort(key=lambda x: x[0])
            melhor = candidatos[0]
            # viés anti-churn: se a config ATUAL cabe nesse pt_b e está quase tão
            # boa quanto a melhor, mantém tudo como está (nenhuma mexida à toa).
            top_atual = limite_inferior_emu - _body_top(int(round(pt_tit_atual))) - FOLGA
            if (abs(pt_b - pt_corpo) < 0.01 and largura0 in altura_larg
                    and altura_larg[largura0] <= top_atual
                    and custo_larg[largura0] <= melhor[0][0] + MARGEM_MASSA):
                plano = (pt_tit_atual, largura0, pt_b)
                break
            plano = (melhor[1], melhor[2], pt_b)
            break
    # nem a 17pt coube: fica no piso, mais largo possível, e avisa (não esmaga)
    if plano is None:
        largura = larguras[-1] if larguras else largura0
        n = _contar_linhas_texto(texto_corpo, int(largura * FATOR_WRAP_RENDER), PT_PISO_CORPO, negrito=False)
        alt = int(n * PT_PISO_CORPO * 12700 * ALTURA_LINHA)
        plano = (float(pt_tit_piso), largura, PT_PISO_CORPO)
        transborda = alt > (limite_inferior_emu - _body_top(pt_tit_piso) - FOLGA)

    pt_tit, largura, pt_b = plano
    # aplica o título no tamanho escolhido (só re-seta se mudou; recalcula viúva +
    # balanceamento nesse pt pra as quebras baterem com a fonte nova)
    if abs(pt_tit - pt_tit_atual) > 0.01:
        tf = _evitar_viuva(texto_titulo, title_box.width, pt_tit)
        tf = _quebrar_titulo_balanceado(tf, title_box.width, pt_tit)
        _set_run_text_keep_format(title_box.text_frame, tf)
        for r in title_box.text_frame.paragraphs[0].runs:
            r.font.size = Pt(pt_tit)
    # aplica a largura e a FONTE do corpo ANTES do texto, e monta a prevenção de
    # viúva já na largura+fonte FINAIS (com _glue_viuva_corpo puro) — senão o nbsp
    # é calculado pra 19pt e aplicado a 18pt, abrindo um vão (bug do "busca").
    if largura != largura0:
        body_box.width = largura
    if abs(pt_b - pt_corpo) > 0.01:
        for paragraph in body_box.text_frame.paragraphs:
            for r in paragraph.runs:
                r.font.size = Pt(pt_b)
    _set_run_text_keep_format(body_box.text_frame, _glue_viuva_corpo(texto_corpo, largura, pt_b, FATOR_WRAP_RENDER))
    if _DEBUG_MASSA:
        print(f"[MASSA] {contexto[:30]!r} L0={largura0/360000:.2f}->{largura/360000:.2f} "
              f"pt_tit={pt_tit:.0f} pt_b={pt_b:.0f} custo={_custo_massa(texto_corpo, largura, pt_b, capa=capa):.3f}")
    # reposiciona o corpo pela altura REAL do título (texto já na caixa)
    _reposicionar_corpo_apos_titulo(title_box, body_box, texto_titulo)
    if transborda and warnings is not None:
        warnings.append(f"{contexto}: muito texto — corpo mantido em {int(pt_b)}pt e "
                        f"ainda pode transbordar; considere encurtar o texto ou ajustar à mão.")


def _ajustar_titulo_largura_reduzida(title_box, texto_titulo, tamanho_min_pt=24):
    """Quando a caixa de título fica mais estreita (ex.: por causa do QR
    code), reduz a fonte até nenhuma palavra ultrapassar a largura —
    evita quebra no meio da palavra e mantém o texto legível."""
    from PIL import ImageFont
    from pptx.util import Pt
    if not title_box.text_frame.paragraphs[0].runs:
        return
    run = title_box.text_frame.paragraphs[0].runs[0]
    tamanho_pt = (run.font.size / 12700) if run.font.size else 40
    tamanho_original = tamanho_pt
    largura_px = (title_box.width / 12700) * 4
    # o espaço inseparável (viúva) gruda duas palavras num "pedaço" só —
    # é esse pedaço que não pode quebrar, então é ele que precisa caber,
    # não a palavra isolada. O \n do rebalanceamento não é pedaço: vira separador.
    pedacos = texto_titulo.replace("\n", " ").split(" ")

    while tamanho_pt > tamanho_min_pt:
        font = ImageFont.truetype(FONT_BOLD, max(int(tamanho_pt * 4), 8))
        maior_largura = max(font.getbbox(p)[2] - font.getbbox(p)[0] for p in pedacos)
        if maior_largura <= largura_px * 0.92:
            break
        tamanho_pt -= 1

    if tamanho_pt != tamanho_original:
        for r in title_box.text_frame.paragraphs[0].runs:
            r.font.size = Pt(tamanho_pt)


# ---------------------------------------------------------------------------
# Utilidades de shape
# ---------------------------------------------------------------------------

def duplicate_slide(prs, template_slide):
    """Cria um slide novo (sempre no fim da apresentação) copiando o XML de
    todos os shapes de `template_slide`, E as relações de imagem — assim
    fotos/tags que SERÃO substituídas funcionam via _replace_picture/_replace_tag,
    e imagens que NÃO são substituídas (logo, estoque de tags não usadas)
    continuam válidas em vez de ficarem quebradas."""
    new_slide = prs.slides.add_slide(template_slide.slide_layout)
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape._element)
    for shape in template_slide.shapes:
        new_el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)
    for rId, rel in template_slide.part.rels.items():
        if rel.reltype == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image":
            new_slide.part.rels._rels[rId] = rel
    return new_slide


def _set_run_text_keep_format(text_frame, new_text):
    para = text_frame.paragraphs[0]
    runs = para.runs
    if not runs:
        return
    runs[0].text = new_text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)
    for p in text_frame.paragraphs[1:]:
        p._p.getparent().remove(p._p)


LARANJA_BOX = (0xFE, 0xA4, 0x2F)  # laranja do rótulo dos boxes (vem do template)


def _preencher_box_textbox(box_shape, label, texto):
    """Preenche uma caixa de box do template: o rótulo mantém o negrito laranja
    que já vem do template e o conteúdo entra em preto regular (igual ao
    modelo do cliente: 'CEA:' laranja + horários em preto)."""
    from pptx.dml.color import RGBColor
    tf = box_shape.text_frame
    para = tf.paragraphs[0]
    runs = para.runs
    if not runs:
        return
    base = runs[0]  # herda formatação (negrito + laranja) do template
    base.text = (label + ": ") if label else ""
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)
    for pextra in tf.paragraphs[1:]:
        pextra._p.getparent().remove(pextra._p)
    run_txt = para.add_run()
    run_txt.text = texto
    run_txt.font.bold = False
    run_txt.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    if base.font.size:
        run_txt.font.size = base.font.size
    if base.font.name:
        run_txt.font.name = base.font.name


def _boxes_no_corpo(body_box, boxes_data):
    """Fallback pra quando o slide NÃO tem placeholders de box (ex.: um layout
    que não seja o do cliente): anexa os boxes ao próprio corpo como parágrafos
    (rótulo laranja em negrito + texto preto). Sem as linhas divisórias, mas o
    conteúdo aparece legível e estruturado em qualquer layout."""
    from pptx.dml.color import RGBColor
    tf = body_box.text_frame
    base_run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
    tam = base_run.font.size if (base_run and base_run.font.size) else None
    for box in boxes_data:
        p = tf.add_paragraph()
        if box["label"]:
            rl = p.add_run()
            rl.text = box["label"] + ": "
            rl.font.bold = True
            rl.font.color.rgb = RGBColor(*LARANJA_BOX)
            if tam:
                rl.font.size = tam
        rt = p.add_run()
        rt.text = box["texto"]
        rt.font.bold = False
        if tam:
            rt.font.size = tam


def _render_boxes(slide, body_box, placeholders, linhas, boxes_data, warnings, contexto=""):
    """Renderiza os boxes da matéria. Estratégia (independente de número de
    slide): se o layout já traz placeholders de box (caixas extras + linhas
    divisórias), PREENCHE em vez de apagar; se não traz, cria os boxes dentro
    do corpo (fallback). Ajusta a quantidade de placeholders/linhas ao número
    real de boxes do briefing."""
    n = len(boxes_data)
    if not placeholders:
        _boxes_no_corpo(body_box, boxes_data)
        return

    # preenche os primeiros placeholders com os boxes disponíveis
    for i in range(min(n, len(placeholders))):
        b = boxes_data[i]
        _preencher_box_textbox(placeholders[i], b["label"], b["texto"])

    # menos boxes que placeholders: remove os que sobraram (e as linhas de baixo)
    sobrando = placeholders[n:]
    for extra in sobrando:
        extra._element.getparent().remove(extra._element)
    if sobrando and linhas:
        for ln in sorted(linhas, key=lambda s: s.top or 0, reverse=True)[:len(sobrando)]:
            ln._element.getparent().remove(ln._element)

    # mais boxes que placeholders: junta o excedente no último box, pra não
    # perder informação, e avisa
    if n > len(placeholders):
        resto = boxes_data[len(placeholders):]
        extra_txt = " ".join(
            ((b["label"] + ": ") if b["label"] else "") + b["texto"] for b in resto)
        p = placeholders[-1].text_frame.add_paragraph()
        r = p.add_run()
        r.text = extra_txt
        warnings.append(
            f"{contexto}: {n} boxes no briefing, mas o layout tem "
            f"{len(placeholders)} — os extras foram juntados no último box.")


def _place_qr_code(slide, qr_path, tag_box, slide_width_emu):
    """Posiciona o QR code (sempre quadrado) encostado na borda direita,
    alinhado verticalmente com a tag da matéria à qual ele pertence."""
    from pptx.util import Emu
    lado = tag_box.height  # QR é quadrado, mesma altura de referência da tag (6.09cm ~ constante)
    lado = Emu(int(6.09 * 360000))
    margem_direita = Emu(int(2.0 * 360000))
    left = slide_width_emu - lado - margem_direita
    top = tag_box.top
    slide.shapes.add_picture(qr_path, left, top, lado, lado)


def _swap_qr_placeholder(slide, qr_path, ignorar=()):
    """Na página única COM QR, o template já traz um placeholder QUADRADO na
    posição exata (definida à mão no PowerPoint). Aqui a gente só troca a
    imagem dele pelo QR real da matéria, mantendo posição e tamanho. Se não
    houver QR pra essa matéria (`qr_path` None), remove o placeholder pra não
    deixar QR de exemplo vazar. Acha o placeholder pela forma (quadrado de
    3–8 cm), ignorando a foto e a tag (`ignorar`). Retorna True se achou."""
    alvo = None
    for s in slide.shapes:
        if s.shape_type != 13 or s.left is None or s.left < 0 or not s.height:
            continue
        if s in ignorar:
            continue
        ratio = s.width / s.height
        lado_cm = s.width / 360000.0
        if 0.85 <= ratio <= 1.18 and 3.0 <= lado_cm <= 8.0:
            alvo = s
            break
    if alvo is None:
        return False
    left, top, w, h = alvo.left, alvo.top, alvo.width, alvo.height
    parent = alvo._element.getparent()
    idx = list(parent).index(alvo._element)
    parent.remove(alvo._element)
    if qr_path:
        new_pic = slide.shapes.add_picture(qr_path, left, top, w, h)
        el = new_pic._element
        el.getparent().remove(el)
        parent.insert(idx, el)   # mantém a ordem-z (QR por cima da foto)
    return True


CORTE_DIAGONAL_PROPORCAO = 0.092  # ~9.2% da altura, medido nas fotos do template

# Viés do corte vertical usado como FALLBACK (quando não há rosto detectado):
# quanto do excedente sai do TOPO (o resto sai da base). 0.5 = centralizado
# (decepa a cabeça de quem está em pé); 0.2 = puxa pro topo, preservando
# cabeça e sacrificando mais o rodapé. Com rosto, quem manda é
# _offset_corte_vertical (dá folga acima da cabeça).
VIES_CORTE_VERTICAL = 0.2


# ---------------------------------------------------------------------------
# Enquadramento inteligente: quando a foto é mais larga que a máscara e precisa
# ser cortada nas laterais, em vez de cortar SEMPRE pelo centro (que decepa
# quem está fora do meio), o motor detecta os rostos e desloca o corte pra
# manter as pessoas enquadradas. Se o OpenCV/modelo não existirem, ou se a foto
# não tiver rosto (prédio, gráfico), cai no corte central de sempre — nunca
# piora e nunca quebra.
# ---------------------------------------------------------------------------
_YUNET = None
_YUNET_TENTOU = False
_ROSTOS_CACHE = {}       # path -> lista de rostos (evita redetecção na mesma build)
_SALIENCIA_CACHE = {}    # path -> (mapa, W0, H0) da saliência


# Reforço dos rostos SOBRE o mapa de saliência. O rosto é assunto forte, mas
# não pode deixar uma MULTIDÃO sequestrar o corte: o peso é pela CONFIANÇA² do
# rosto (o palestrante nítido/frontal, score alto, pesa muito mais que as nucas
# da plateia) vezes a raiz da área (a área conta, mas sem explodir). Peso baixo
# de propósito — a saliência é quem manda; o rosto só ajusta.
FATOR_REFORCO_ROSTO = 0.004
# Acima de tantos rostos a cena é "multidão": o snap de borda é desligado (senão
# trocaria o assunto por um rosto qualquer encostado na borda da janela).
MAX_ROSTOS_PARA_SNAP = 3


def _get_yunet():
    """Carrega o detector de rosto YuNet uma vez só (cacheado). Devolve None se
    o OpenCV ou o modelo não estiverem disponíveis."""
    global _YUNET, _YUNET_TENTOU
    if _YUNET_TENTOU:
        return _YUNET
    _YUNET_TENTOU = True
    try:
        import cv2
        # silencia o ruído do OpenCV no terminal (ex.: o WARN
        # "setPreferableTarget Targets are not supported by the new graph
        # engine" do OpenCV 5 ao rodar o YuNet) — é só aviso, não erro; o
        # enquadramento funciona igual. Mantém o nível ERROR (erros reais ainda
        # aparecem). Protegido: se a API de log mudar, não quebra.
        try:
            cv2.utils.logging.setLogLevel(cv2.utils.logging.LOG_LEVEL_ERROR)
        except Exception:
            pass
        if _os.path.isfile(YUNET_PATH):
            _YUNET = cv2.FaceDetectorYN.create(YUNET_PATH, "", (320, 320),
                                               score_threshold=0.6)
    except Exception:
        _YUNET = None
    return _YUNET


def _detectar_rostos(image_path):
    """Rostos (x, y, w, h, score) em coordenadas da imagem original, cacheado
    por caminho (a mesma foto é consultada várias vezes numa build)."""
    if image_path not in _ROSTOS_CACHE:
        _ROSTOS_CACHE[image_path] = _detectar_rostos_raw(image_path)
    return _ROSTOS_CACHE[image_path]


def _detectar_rostos_raw(image_path):
    """Rostos (x, y, w, h, score) em coordenadas da imagem original. Roda numa
    versão reduzida (~1024px) pra velocidade. Devolve [] se não houver
    detector/modelo/rosto — aí o chamador usa o enquadramento padrão."""
    det = _get_yunet()
    if det is None:
        return []
    try:
        import cv2
        import numpy as np
        # carrega via PIL (como o resto do motor) e converte pro formato do
        # OpenCV — evita o "libpng warning: iCCP" que o cv2.imread cospe no
        # terminal ao ler PNGs com perfil de cor levemente torto (é só ruído,
        # mas suja a saída). PIL lê essas PNGs sem reclamar.
        img = cv2.cvtColor(np.array(Image.open(image_path).convert("RGB")),
                           cv2.COLOR_RGB2BGR)
        H, W = img.shape[:2]
        escala = 1024.0 / W if W > 1024 else 1.0
        sw, sh = max(int(W * escala), 1), max(int(H * escala), 1)
        det.setInputSize((sw, sh))
        entrada = cv2.resize(img, (sw, sh)) if escala != 1.0 else img
        _, faces = det.detect(entrada)
        if faces is None:
            return []
        return [(float(f[0] / escala), float(f[1] / escala),
                 float(f[2] / escala), float(f[3] / escala), float(f[-1]))
                for f in faces]
    except Exception:
        return []


# ---------------------------------------------------------------------------
# Saliência visual (numpy puro, SEM dependência nova): responde "onde está o
# assunto?" quando o rosto sozinho não basta. Método do resíduo espectral
# (Hou & Zhang): realça o que é SINGULAR na imagem (o palestrante isolado, o
# globo brilhante) e SUPRIME padrões repetitivos (a textura de cadeiras/nucas
# de uma plateia). É isso que impede a multidão de sequestrar o corte e o que
# reposiciona uma foto sem rosto cujo assunto está fora do centro.
# ---------------------------------------------------------------------------
def _box_blur(a, k):
    """Média de janela k x k via soma acumulada (numpy puro, sem scipy)."""
    import numpy as np
    if k < 2:
        return a
    pad = k // 2
    ap = np.pad(a, pad, mode="edge")
    cs = np.cumsum(np.cumsum(ap, axis=0), axis=1)
    cs = np.pad(cs, ((1, 0), (1, 0)), mode="constant")
    H, W = a.shape
    y1 = np.arange(H)[:, None]; x1 = np.arange(W)[None, :]
    y2 = y1 + k; x2 = x1 + k
    tot = cs[y2, x2] - cs[y1, x2] - cs[y2, x1] + cs[y1, x1]
    return tot / (k * k)


def _mapa_saliencia(image_path, largura=64):
    """Mapa de saliência (0..1) por resíduo espectral, numa versão minúscula da
    foto. Devolve (mapa, W0, H0) ou None se numpy/Pillow falharem — aí o
    chamador cai no corte de sempre (nunca quebra, nunca piora)."""
    if image_path in _SALIENCIA_CACHE:
        return _SALIENCIA_CACHE[image_path]
    dados = None
    try:
        import numpy as np
        im = Image.open(image_path).convert("L")
        W0, H0 = im.size
        h = max(int(round(largura * H0 / W0)), 16)
        g = np.asarray(im.resize((largura, h), Image.BILINEAR), dtype=np.float64)
        F = np.fft.fft2(g)
        logamp = np.log(np.abs(F) + 1e-8)
        fase = np.angle(F)
        residuo = logamp - _box_blur(logamp, 3)
        S = np.abs(np.fft.ifft2(np.exp(residuo + 1j * fase))) ** 2
        S = _box_blur(S, 5)
        S -= S.min()
        if S.max() > 0:
            S /= S.max()
        dados = (S, W0, H0)
    except Exception:
        dados = None
    _SALIENCIA_CACHE[image_path] = dados
    return dados


def _centro_importancia(image_path, rostos, eixo):
    """Centro do 'assunto' ao longo de um eixo (0 = x/horizontal, 1 = y/
    vertical): energia da saliência + reforço dos rostos. Devolve px na imagem
    original, ou None se não houver sinal (aí volta pro corte padrão)."""
    dados = _mapa_saliencia(image_path)
    if dados is None:
        return None
    try:
        import numpy as np
        S, W0, H0 = dados
        S = S.copy()
        hs, ws = S.shape
        esc_x = W0 / ws
        esc_y = H0 / hs
        for (x, y, w, h, s) in rostos:
            gx0 = max(0, int(x / esc_x)); gx1 = min(ws, max(int((x + w) / esc_x), gx0 + 1))
            gy0 = max(0, int(y / esc_y)); gy1 = min(hs, max(int((y + h) / esc_y), gy0 + 1))
            S[gy0:gy1, gx0:gx1] += (s ** 2) * np.sqrt(w * h) * FATOR_REFORCO_ROSTO
        if eixo == 0:
            perfil = S.sum(axis=0); esc = esc_x
        else:
            perfil = S.sum(axis=1); esc = esc_y
        total = float(perfil.sum())
        if total <= 0:
            return None
        idx = np.arange(len(perfil))
        centro = float((idx * perfil).sum() / total)
        return (centro + 0.5) * esc
    except Exception:
        return None


# O corte central só é abandonado se mover pro assunto capturar mais que este
# tanto da energia de saliência. Protege logo/gráfico já bem centrado de ser
# descentralizado por uma saliência fraca e levemente assimétrica (ex.: o peso
# de tinta de um texto puxando o centro-de-massa pro lado).
LIMIAR_ENERGIA_DESCENTRAR = 0.95


def _offset_horizontal_saliencia_sem_rosto(image_path, largura_total, new_w):
    """x0 do corte horizontal para foto SEM rosto: usa a saliência pra achar o
    assunto (globo brilhante, objeto fora do centro), mas só sai do centro se
    a janela no assunto capturar materialmente mais energia que a janela
    central — senão devolve o centro (não mexe no que já estava bem centrado).
    Devolve None se a saliência não estiver disponível."""
    dados = _mapa_saliencia(image_path)
    if dados is None:
        return None
    try:
        import numpy as np
        S, W0, H0 = dados
        col = S.sum(axis=0)
        ws = len(col)
        esc = W0 / ws
        total = float(col.sum())
        if total <= 0:
            return None
        lo, hi = 0, largura_total - new_w
        x0_mid = (largura_total - new_w) // 2
        centro = (float((np.arange(ws) * col).sum() / total) + 0.5) * esc
        x0_cent = max(lo, min(hi, int(round(centro - new_w / 2.0))))

        def energia(x0px):
            g0 = x0px / esc
            g1 = (x0px + new_w) / esc
            i0 = max(0, int(np.floor(g0)))
            i1 = min(ws, int(np.ceil(g1)))
            return float(col[i0:i1].sum())

        e_assunto = energia(x0_cent)
        if e_assunto <= 0:
            return x0_mid
        if energia(x0_mid) >= LIMIAR_ENERGIA_DESCENTRAR * e_assunto:
            return x0_mid                 # o centro já pega quase tudo: não descentraliza
        return x0_cent
    except Exception:
        return None


def _parse_foco(foco):
    """Interpreta o 'FOCO:' do briefing -> (direcao, fracao).
    - direcao: 'esquerda'/'direita'/'centro'/'cima'/'baixo' (ou None).
    - fracao: 0.0..1.0 = o quanto empurrar A PARTIR do enquadramento automático
      rumo àquela borda. None quando não veio número -> vale o comportamento
      antigo (empurra até a borda = 100%). Ex.: 'direita 20' -> ('direita', 0.2);
      'direita' -> ('direita', None); 'centro' -> ('centro', None)."""
    if not foco:
        return None, None
    s = unicodedata.normalize("NFKD", str(foco)).encode("ascii", "ignore").decode().lower()
    direcao = None
    for chave, canon in (
        ("esquerda", "esquerda"), ("esq", "esquerda"), ("left", "esquerda"),
        ("direita", "direita"), ("dir", "direita"), ("right", "direita"),
        ("centro", "centro"), ("center", "centro"), ("meio", "centro"),
        ("cima", "cima"), ("topo", "cima"), ("top", "cima"),
        ("baixo", "baixo"), ("base", "baixo"), ("bottom", "baixo"), ("fundo", "baixo"),
    ):
        if chave in s:
            direcao = canon
            break
    frac = None
    m = re.search(r"(\d+(?:[.,]\d+)?)", s)
    if m:
        frac = max(0.0, min(1.0, float(m.group(1).replace(",", ".")) / 100.0))
    return direcao, frac


# palavras-chave de cada eixo, da MAIS LONGA pra mais curta (pra 'direita'
# ganhar de 'dir' e 'esquerda' de 'esq' na alternância do regex).
_FOCO_KEYS_H = [("esquerda", "esquerda"), ("left", "esquerda"), ("esq", "esquerda"),
                ("direita", "direita"), ("right", "direita"), ("dir", "direita"),
                ("centro", "centro"), ("center", "centro"), ("meio", "centro")]
_FOCO_KEYS_V = [("cima", "cima"), ("topo", "cima"), ("top", "cima"),
                ("baixo", "baixo"), ("base", "baixo"), ("bottom", "baixo"),
                ("fundo", "baixo"), ("centro", "centro"), ("center", "centro"),
                ("meio", "centro")]


def _parse_foco_eixo(foco, keys):
    """(direcao_canon, fracao) para UM eixo. Acha a 1ª palavra-chave DAQUELE eixo
    na string e o número que vem logo depois dela — assim um foco 2D
    'direita 30 baixo 20' extrai 30 no horizontal e 20 no vertical (cada eixo
    lê o número que segue a SUA palavra, não o 1º número da linha)."""
    if not foco:
        return None, None
    s = unicodedata.normalize("NFKD", str(foco)).encode("ascii", "ignore").decode().lower()
    alt = "|".join(re.escape(k) for k, _ in keys)
    m = re.search(r"(" + alt + r")\s*(\d+(?:[.,]\d+)?)?", s)
    if not m:
        return None, None
    canon = dict(keys)[m.group(1)]
    frac = None
    if m.group(2):
        frac = max(0.0, min(1.0, float(m.group(2).replace(",", ".")) / 100.0))
    return canon, frac


def _foco_horizontal(foco):
    """(direcao, fracao) do eixo HORIZONTAL (esquerda/direita/centro)."""
    return _parse_foco_eixo(foco, _FOCO_KEYS_H)


def _foco_vertical(foco):
    """(direcao, fracao) do eixo VERTICAL (cima/baixo/centro)."""
    return _parse_foco_eixo(foco, _FOCO_KEYS_V)


def _empurrar_para_borda(auto, lo, hi, direcao_baixa, direcao_alta, direcao, frac):
    """Aplica o override direcional sobre o offset automático `auto`: empurra
    `frac` do caminho entre `auto` e a borda (lo ou hi). frac None = borda (100%,
    o comportamento antigo do FOCO sem número)."""
    if direcao == direcao_baixa:
        f = 1.0 if frac is None else frac
        return int(round(max(lo, min(hi, auto - f * (auto - lo)))))
    if direcao == direcao_alta:
        f = 1.0 if frac is None else frac
        return int(round(max(lo, min(hi, auto + f * (hi - auto)))))
    return auto  # centro / None -> mantém o automático


def _centro_horizontal_rostos(rostos):
    """Centro x ponderado pela ÁREA de cada rosto: os rostos maiores/mais
    próximos (os assuntos principais) pesam mais, e um eventual falso-positivo
    pequeno e distante quase não desloca o corte."""
    soma = peso_total = 0.0
    for (x, y, w, h, s) in rostos:
        peso = w * h
        soma += (x + w / 2.0) * peso
        peso_total += peso
    return (soma / peso_total) if peso_total > 0 else None


def _offset_corte_horizontal(new_image_path, largura_total, new_w, foco=None):
    """x0 do corte horizontal. Sem FOCO (ou 'FOCO: centro') usa o enquadramento
    automático (assunto). 'FOCO: esquerda/direita' empurra a PARTIR do
    automático rumo à borda; com porcentagem ('direita 20') empurra só essa
    fração do caminho, sem número vai até a borda (o comportamento antigo)."""
    lo, hi = 0, largura_total - new_w
    auto = _offset_horizontal_auto(new_image_path, largura_total, new_w)
    direcao, frac = _foco_horizontal(foco)
    # convenção: o ASSUNTO se move na direção da seta. "direita" (assunto vai pra
    # direita do quadro) = janela pra ESQUERDA = x0 menor (lo). "esquerda" =
    # janela pra direita = x0 maior (hi).
    return _empurrar_para_borda(auto, lo, hi, "direita", "esquerda", direcao, frac)


def _offset_horizontal_auto(new_image_path, largura_total, new_w):
    """Enquadramento horizontal AUTOMÁTICO (sem override), pelo ASSUNTO:
    1) centro de importância = saliência (suprime multidão, acha o singular) +
       reforço dos rostos;
    2) fallback ao centro ponderado dos rostos / centro geométrico se a
       saliência não estiver disponível.
    Quando TODOS os rostos cabem na janela, ainda garante contê-los (mantém o
    comportamento de sempre pras fotos de 1–3 pessoas — baixa regressão)."""
    lo, hi = 0, largura_total - new_w
    x0_centro = (largura_total - new_w) // 2

    rostos = _detectar_rostos(new_image_path)
    if not rostos:
        # sem rosto: saliência com guarda de energia (não descentraliza gráfico
        # já centrado); cai no centro se a saliência não estiver disponível
        x0 = _offset_horizontal_saliencia_sem_rosto(new_image_path, largura_total, new_w)
        return x0 if x0 is not None else x0_centro

    rx0 = min(r[0] for r in rostos)
    rx1 = max(r[0] + r[2] for r in rostos)
    todos_cabem = (rx1 - rx0) <= new_w
    rosto_alvo = None            # rosto único priorizado (quando poucos não cabem)
    if todos_cabem:
        # as pessoas cabem inteiras na janela: o assunto são ELAS → centraliza
        # no centro dos rostos. A saliência NÃO entra aqui de propósito: num
        # trio contra um fundo laranja/luminoso, a saliência do fundo puxava o
        # corte e jogava as pessoas pro lado. Com rosto que cabe, quem manda é
        # o rosto (comportamento estável de sempre).
        centro = _centro_horizontal_rostos(rostos)
        if centro is None:
            centro = x0_centro + new_w / 2.0
    elif len(rostos) <= MAX_ROSTOS_PARA_SNAP:
        # POUCAS pessoas que não cabem juntas na janela (ex.: um casal numa fatia
        # vertical estreita de uma união): em vez de centralizar no VÃO entre elas
        # — que corta os dois rostos —, prioriza enquadrar UMA, a mais proeminente
        # (área x confiança). É a decisão do olho: se não cabem os dois, mostra
        # um bem enquadrado. (Ex.: CEDROS BIER 3, casal lado a lado.)
        rosto_alvo = max(rostos, key=lambda r: r[2] * r[3] * (r[4] if len(r) > 4 else 1.0))
        centro = rosto_alvo[0] + rosto_alvo[2] / 2.0
    else:
        # multidão (muitos rostos que não cabem): aí a saliência acha o assunto
        # singular e suprime a repetição da plateia.
        centro = _centro_importancia(new_image_path, rostos, eixo=0)
        if centro is None:
            c = _centro_horizontal_rostos(rostos)
            centro = c if c is not None else (x0_centro + new_w / 2.0)
    x0 = int(round(centro - new_w / 2.0))
    x0 = max(lo, min(hi, x0))
    if todos_cabem:                                # garante conter todos os rostos
        x0 = min(max(x0, int(rx1 - new_w)), int(rx0))
        x0 = max(lo, min(hi, x0))
    elif rosto_alvo is not None:                   # contém o rosto priorizado inteiro
        fx0, fw = rosto_alvo[0], rosto_alvo[2]
        x0 = min(max(x0, int(fx0 + fw - new_w)), int(fx0))
        x0 = max(lo, min(hi, x0))
    # snap só com POUCOS rostos (sujeito claro). Em multidão, o snap trocaria
    # o assunto por um rosto qualquer de borda — então é desligado.
    if len(rostos) <= MAX_ROSTOS_PARA_SNAP:
        x0 = _snap_fora_dos_rostos(x0, new_w, rostos, lo, hi)
    return x0


def _snap_fora_dos_rostos(x0, new_w, rostos, lo, hi):
    def fatia(pos):
        return any(x < pos < x + w or x < pos + new_w < x + w
                   for (x, y, w, h, s) in rostos)
    if not fatia(x0):
        return x0
    candidatos = [lo, hi]
    for (x, y, w, h, s) in rostos:
        # posições que encostam a borda ESQ ou DIR da janela na borda de um rosto
        candidatos += [x, x + w, x - new_w, x + w - new_w]
    validos = [c for c in candidatos if lo <= c <= hi and not fatia(c)]
    return int(min(validos, key=lambda c: abs(c - x0))) if validos else x0


def _offset_corte_vertical(new_image_path, altura_total, new_h, foco=None):
    """y0 do corte vertical. Sem FOCO (ou 'FOCO: centro') usa o automático.
    'FOCO: cima/baixo' empurra a PARTIR do automático rumo à borda; com
    porcentagem ('baixo 20') empurra só essa fração, sem número vai até a
    borda."""
    lo, hi = 0, altura_total - new_h
    auto = _offset_vertical_auto(new_image_path, altura_total, new_h)
    direcao, frac = _foco_vertical(foco)
    # mesma convenção do horizontal: o assunto se move na direção da seta.
    # "baixo" (assunto pra baixo) = janela pra CIMA = y0 menor (lo); "cima" =
    # janela pra baixo = y0 maior (hi).
    return _empurrar_para_borda(auto, lo, hi, "baixo", "cima", direcao, frac)


def _offset_vertical_auto(new_image_path, altura_total, new_h):
    """Corte vertical AUTOMÁTICO: com rostos, dá FOLGA acima da cabeça do rosto
    mais alto (não decepa cabeça/cabelo/boné) e joga o excedente pra BASE; sem
    rosto, usa o viés fixo pro topo (VIES_CORTE_VERTICAL), o comportamento
    antigo — o vertical automático fica intocado."""
    lo, hi = 0, altura_total - new_h
    y0_padrao = int((altura_total - new_h) * VIES_CORTE_VERTICAL)
    rostos = _detectar_rostos(new_image_path)
    if not rostos:
        return y0_padrao
    fy0 = min(y for (x, y, w, h, s) in rostos)                 # topo do rosto mais alto
    fh = next(h for (x, y, w, h, s) in rostos if y == fy0)     # altura desse rosto
    folga = int(0.7 * fh)                                       # cabelo/cabeça/boné acima da caixa
    y0 = max(lo, min(hi, fy0 - folga))
    # se a BASE da janela cortar um rosto ao meio, SOBE só o necessário pra
    # deixá-lo inteiro fora (nunca desce, pra não voltar a cortar as cabeças)
    for (x, y, w, h, s) in sorted(rostos, key=lambda r: r[1]):
        if y < y0 + new_h < y + h:
            y0 = max(lo, y - new_h)
            break
    return y0


# zoom aplicado SÓ quando o briefing traz foco nos DOIS eixos (opt-in). O corte
# "cover" normal só deixa folga num eixo; um zoom modesto cria folga nos dois,
# permitindo o ajuste 2D. Com um foco só (ou nenhum) isto nunca é usado e o
# corte é byte-a-byte o de sempre.
ZOOM_FOCO_2D = 1.18


def _crop_zoom_2d(new_image_path, target_ratio, foco, zoom=ZOOM_FOCO_2D):
    """Corte com foco nos DOIS eixos ao mesmo tempo. Dá um zoom modesto pra
    gerar folga horizontal E vertical (o cover só sobra num eixo) e posiciona
    cada eixo pelo seu foco, a PARTIR do enquadramento automático do assunto.
    Só é chamado quando há foco horizontal + vertical juntos."""
    im = Image.open(new_image_path).convert("RGB")
    w, h = im.size
    src_ratio = w / h
    if src_ratio > target_ratio:            # foto larga: cover pela altura
        cover_w, cover_h = h * target_ratio, h
    else:                                   # foto alta: cover pela largura
        cover_w, cover_h = w, w / target_ratio
    crop_w = min(w, int(round(cover_w / zoom)))
    crop_h = min(h, int(round(cover_h / zoom)))

    auto_x = _offset_horizontal_auto(new_image_path, w, crop_w)
    dh, fh = _foco_horizontal(foco)
    x0 = _empurrar_para_borda(auto_x, 0, w - crop_w, "direita", "esquerda", dh, fh)

    auto_y = _offset_vertical_auto(new_image_path, h, crop_h)
    dv, fv = _foco_vertical(foco)
    y0 = _empurrar_para_borda(auto_y, 0, h - crop_h, "baixo", "cima", dv, fv)

    return im.crop((x0, y0, x0 + crop_w, y0 + crop_h))


def _aplicar_corte_diagonal(im, proporcao=CORTE_DIAGONAL_PROPORCAO, cor_fundo=(255, 255, 255)):
    """Pinta o triângulo no canto superior-esquerdo, reproduzindo o corte
    diagonal que já vem queimado nas fotos do Jornal Mural."""
    im = im.copy()
    w, h = im.size
    perna = int(round(h * proporcao))
    if perna <= 0:
        return im
    px = im.load()
    for y in range(perna):
        limite_x = perna - y  # reta de (perna,0) a (0,perna)
        for x in range(min(limite_x, w)):
            px[x, y] = cor_fundo
    return im


def _compose_blurred_fit(image_path, largura_alvo_px, altura_alvo_px, escurecer=0.7):
    """Encaixa a foto INTEIRA na caixa (sem cortar nada), preenchendo o vazio
    com uma versão ampliada e desfocada da própria foto. É o acabamento que
    se faz à mão pra fotos que não podem ser cortadas — banner, telão, arte
    com texto nas bordas: em vez de perder as laterais, mostra tudo e o fundo
    borrado dá o preenchimento, sem barra preta dura."""
    from PIL import ImageFilter, ImageEnhance
    im = Image.open(image_path).convert("RGB")
    alvo_ratio = largura_alvo_px / altura_alvo_px
    w, h = im.size
    src_ratio = w / h

    # fundo: a foto preenchendo toda a caixa (cover), bem desfocada e um
    # pouco escurecida, virando uma textura neutra atrás da peça
    if src_ratio > alvo_ratio:
        nb = int(h * alvo_ratio)
        x0 = (w - nb) // 2
        fundo = im.crop((x0, 0, x0 + nb, h))
    else:
        nb = int(w / alvo_ratio)
        y0 = (h - nb) // 2
        fundo = im.crop((0, y0, w, y0 + nb))
    fundo = fundo.resize((largura_alvo_px, altura_alvo_px))
    fundo = fundo.filter(ImageFilter.GaussianBlur(
        radius=max(largura_alvo_px, altura_alvo_px) // 30))
    fundo = ImageEnhance.Brightness(fundo).enhance(escurecer)

    # frente: a foto inteira (contain), centralizada sobre o fundo
    escala = min(largura_alvo_px / w, altura_alvo_px / h)
    fw, fh = int(w * escala), int(h * escala)
    frente = im.resize((fw, fh))
    canvas = fundo.copy()
    canvas.paste(frente, ((largura_alvo_px - fw) // 2, (altura_alvo_px - fh) // 2))
    return canvas


def _replace_picture(picture_shape, new_image_path, corte_diagonal=True,
                     foto_inteira=False, foco=None, warnings=None, contexto=""):
    left, top, width, height = (picture_shape.left, picture_shape.top,
                                 picture_shape.width, picture_shape.height)
    target_ratio = width / height

    if foto_inteira:
        # acabamento "foto inteira": mostra tudo, com fundo desfocado. Renderiza
        # o composto na proporção exata da caixa, a ~150 dpi.
        dpi = 150
        tw = max(int(width / 914400 * dpi), 300)
        th = max(int(height / 914400 * dpi), 300)
        im = _compose_blurred_fit(new_image_path, tw, th)
    else:
        im = Image.open(new_image_path).convert("RGB")
        w, h = im.size
        src_ratio = w / h
        dh, _fh = _foco_horizontal(foco)
        dv, _fv = _foco_vertical(foco)
        dois_eixos = dh in ("esquerda", "direita") and dv in ("cima", "baixo")
        if dois_eixos:
            # foco 2D (opt-in): há foco horizontal E vertical → zoom modesto pra
            # dar folga nos dois eixos e posicionar nos dois. Um foco só cai no
            # ramo de baixo, com corte byte-a-byte idêntico ao de sempre.
            im = _crop_zoom_2d(new_image_path, target_ratio, foco)
        else:
            # avisa se o FOCO pedido é pro eixo que essa foto NÃO corta (ex.:
            # pediu cima/baixo mas a foto é cortada nas laterais) — senão o
            # override fica "sem efeito" calado, o que confunde.
            if foco and warnings is not None:
                _dir, _ = _parse_foco(foco)
                horizontal = src_ratio > target_ratio
                if _dir in ("esquerda", "direita") and not horizontal:
                    warnings.append(f"{contexto}: FOCO '{foco}' (lateral) sem efeito — essa "
                                    f"foto é cortada em cima/baixo; use cima/baixo.")
                elif _dir in ("cima", "baixo") and horizontal:
                    warnings.append(f"{contexto}: FOCO '{foco}' (cima/baixo) sem efeito — essa "
                                    f"foto é cortada nas laterais; use esquerda/direita.")
            if src_ratio > target_ratio:
                new_w = int(h * target_ratio)
                # enquadra pelo ASSUNTO (saliência + rostos); respeita o override
                # FOCO: do briefing; cai no centro se não houver sinal
                x0 = _offset_corte_horizontal(new_image_path, w, new_w, foco=foco)
                im = im.crop((x0, 0, x0 + new_w, h))
            else:
                new_h = int(w / target_ratio)
                # dá folga acima das cabeças e joga o excedente pra base (cai no
                # viés fixo pro topo se não houver rosto); respeita o FOCO: vertical
                y0 = _offset_corte_vertical(new_image_path, h, new_h, foco=foco)
                im = im.crop((0, y0, w, y0 + new_h))

    if corte_diagonal:
        im = _aplicar_corte_diagonal(im)

    buf = io.BytesIO()
    im.save(buf, format="JPEG", quality=92)
    buf.seek(0)

    slide = picture_shape.part.slide
    old_el = picture_shape._element
    sp_parent = old_el.getparent()
    idx = list(sp_parent).index(old_el)
    sp_parent.remove(old_el)
    new_pic = slide.shapes.add_picture(buf, left, top, width, height)
    new_el = new_pic._element
    new_el.getparent().remove(new_el)
    sp_parent.insert(idx, new_el)


def _resolver_editoria(editoria_name):
    key = editoria_name.strip().upper()
    if key in EDITORIA_TAGS:
        return key
    # tenta casar por prefixo/substring (ex.: "EXCELÊNCIA" -> "EXCELÊNCIA OPERACIONAL")
    candidatos = [k for k in EDITORIA_TAGS if k.startswith(key) or key in k]
    if len(candidatos) == 1:
        return candidatos[0]
    return None


def _replace_tag(tag_shape, editoria_name, limite_direito_emu=None):
    """Troca a tag de editoria pela arte certa. A tag mantém a altura do
    template e a largura sai da proporção do PNG — então editoria de nome
    comprido ('INOVAÇÃO E TECNOLOGIA') gera uma tag mais larga. Se essa
    largura fizer a tag passar do `limite_direito_emu` (ex.: a borda esquerda
    da foto na capa), a tag é reduzida proporcionalmente (largura E altura)
    até caber, em vez de invadir a foto."""
    from pptx.util import Emu
    key = _resolver_editoria(editoria_name)
    fname = EDITORIA_TAGS.get(key) if key else None
    if not fname:
        raise ValueError(f"Editoria '{editoria_name}' não mapeada em EDITORIA_TAGS. "
                          f"Adicione o par nome->arquivo antes de gerar.")
    path = _os.path.join(TAGS_DIR, fname)
    left, top, height = tag_shape.left, tag_shape.top, tag_shape.height
    im = Image.open(path)
    w, h = im.size
    new_width = int(height * (w / h))

    slide = tag_shape.part.slide

    # nunca deixa a tag passar da borda direita do slide; e, se veio um limite
    # mais apertado (a foto, na capa), respeita esse limite.
    pres = slide.part.package.presentation_part.presentation
    limite = limite_direito_emu if limite_direito_emu is not None \
        else pres.slide_width - Emu(int(0.5 * 360000))
    if left + new_width > limite > left:
        escala = (limite - left) / new_width
        new_width = int(new_width * escala)
        height = int(height * escala)

    old_el = tag_shape._element
    sp_parent = old_el.getparent()
    idx = list(sp_parent).index(old_el)
    sp_parent.remove(old_el)
    new_pic = slide.shapes.add_picture(path, left, top, new_width, height)
    new_el = new_pic._element
    new_el.getparent().remove(new_el)
    sp_parent.insert(idx, new_el)
    return new_pic


def _top_cm(shape):
    return shape.top / 360000


# ---------------------------------------------------------------------------
# Casamento de arquivo de foto (nome do briefing -> arquivo real na pasta)
# ---------------------------------------------------------------------------

def _normalize(s):
    s = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode()
    return re.sub(r"[^a-z0-9]", "", s.lower())


def _compose_merged_photo(photo_paths, target_ratio, focos=None):
    """Junta 2 ou mais fotos lado a lado (com um filete branco de separação),
    cada uma enquadrada pelo ASSUNTO (mesmo corte das fotos normais: saliência +
    rostos, evita cortar alguém ao meio; cai no centro se não houver sinal),
    formando UMA imagem só na proporção da caixa. `focos` é uma lista alinhada a
    `photo_paths` com o override de enquadramento POR FOTO (esquerda/direita/
    cima/baixo/%); None numa posição = enquadramento automático daquela foto. O
    corte diagonal do canto é aplicado depois, sobre o conjunto inteiro."""
    n = len(photo_paths)
    if focos is None:
        focos = [None] * n
    divisor_frac = 0.008
    altura_ref = 1000
    largura_total = int(altura_ref * target_ratio)
    divisor_px = max(2, int(largura_total * divisor_frac))
    largura_util = largura_total - divisor_px * (n - 1)
    largura_cada = largura_util // n

    canvas = Image.new("RGB", (largura_total, altura_ref), (255, 255, 255))
    x = 0
    sub_ratio = largura_cada / altura_ref
    for path, foco in zip(photo_paths, focos):
        im = Image.open(path).convert("RGB")
        w, h = im.size
        src_ratio = w / h
        if src_ratio > sub_ratio:
            new_w = int(h * sub_ratio)
            # enquadra pelo assunto (saliência + rostos); respeita o FOCO POR
            # FOTO da união (None = automático)
            x0 = _offset_corte_horizontal(path, w, new_w, foco=foco)
            im = im.crop((x0, 0, x0 + new_w, h))
        else:
            new_h = int(w / sub_ratio)
            y0 = _offset_corte_vertical(path, h, new_h, foco=foco)  # folga acima das cabeças
            im = im.crop((0, y0, w, y0 + new_h))
        im = im.resize((largura_cada, altura_ref))
        canvas.paste(im, (x, 0))
        x += largura_cada + divisor_px

    return canvas


def _mosaic_layouts(n):
    """Layouts candidatos do mosaico por nº de fotos. Cada célula é
    (x, y, w, h) em FRAÇÃO da caixa. Há mais de um candidato por nº quando faz
    sentido (ex.: tudo deitado pede linhas largas; tudo em pé, colunas altas) —
    o compositor escolhe o que melhor casa com as orientações reais das fotos."""
    if n == 2:
        return [
            [(0, 0, 0.5, 1.0), (0.5, 0, 0.5, 1.0)],          # 2 colunas altas
            [(0, 0, 1.0, 0.5), (0, 0.5, 1.0, 0.5)],          # 2 faixas largas
        ]
    if n == 3:
        return [
            # 1 alta à esquerda + 2 à direita (o layout do David: a foto deitada
            # ganha uma célula larga em cima e cabe inteira)
            [(0, 0, 0.5, 1.0), (0.5, 0, 0.5, 0.5), (0.5, 0.5, 0.5, 0.5)],
            [(0.5, 0, 0.5, 1.0), (0, 0, 0.5, 0.5), (0, 0.5, 0.5, 0.5)],   # espelhado
            [(0, 0, 1.0, 1/3), (0, 1/3, 1.0, 1/3), (0, 2/3, 1.0, 1/3)],   # 3 faixas largas
            [(0, 0, 1/3, 1.0), (1/3, 0, 1/3, 1.0), (2/3, 0, 1/3, 1.0)],   # 3 colunas altas
        ]
    if n == 4:
        return [
            [(0, 0, 0.5, 0.5), (0.5, 0, 0.5, 0.5),
             (0, 0.5, 0.5, 0.5), (0.5, 0.5, 0.5, 0.5)],                   # 2x2
            [(0, 0, 0.5, 1.0), (0.5, 0, 0.5, 1/3),
             (0.5, 1/3, 0.5, 1/3), (0.5, 2/3, 0.5, 1/3)],                 # 1 alta + 3
        ]
    # n>=5 (raro): grade em faixas iguais de colunas
    return [[(i / n, 0, 1 / n, 1.0) for i in range(n)]]


def _compose_mosaic_photo(photo_paths, target_ratio, focos=None):
    """Compõe as fotos num MOSAICO de células de tamanhos/orientações variados
    (opção '(mosaico)' do briefing), em vez de fatias verticais iguais. Serve pra
    resguardar fotos que NÃO cabem num corte vertical estreito (deitadas, com 2
    pessoas ou produto): elas ganham uma célula larga e aparecem inteiras.

    Melhoria sobre a máscara fixa: o motor casa a ORIENTAÇÃO de cada foto com a
    célula de proporção parecida (deitada→larga, em pé→alta), escolhendo, entre
    os layouts candidatos, o que corta menos as fotos daquela matéria. Cada
    célula é recortada com o mesmo enquadramento inteligente (rosto/saliência) e
    respeita o FOCO por foto. Sai UMA imagem só, do tamanho da caixa."""
    import itertools, math
    n = len(photo_paths)
    if n < 2:
        return _compose_merged_photo(photo_paths, target_ratio, focos)
    if focos is None:
        focos = [None] * n
    altura_ref = 1000
    largura_total = int(altura_ref * target_ratio)
    divisor_px = max(2, int(largura_total * 0.008))

    ratios = []
    for p in photo_paths:
        with Image.open(p) as im:
            w, h = im.size
        ratios.append(w / h)

    # escolhe (layout, atribuição foto→célula) que MINIMIZA o desencaixe de
    # orientação: soma de |log(ratio_foto) - log(ratio_célula)| (log = trata
    # "2x mais largo" e "2x mais alto" como desencaixes iguais).
    melhor = None
    for layout in _mosaic_layouts(n):
        cell_ratios = [(wf * largura_total) / (hf * altura_ref) for (_x, _y, wf, hf) in layout]
        for perm in itertools.permutations(range(n)):   # n<=4 na prática: barato
            custo = sum(abs(math.log(ratios[i]) - math.log(cell_ratios[perm[i]])) for i in range(n))
            if melhor is None or custo < melhor[0]:
                melhor = (custo, layout, perm)
    _, layout, perm = melhor

    canvas = Image.new("RGB", (largura_total, altura_ref), (255, 255, 255))
    m = divisor_px // 2
    for i, path in enumerate(photo_paths):
        xf, yf, wf, hf = layout[perm[i]]
        x0, y0 = int(round(xf * largura_total)), int(round(yf * altura_ref))
        x1, y1 = int(round((xf + wf) * largura_total)), int(round((yf + hf) * altura_ref))
        # filete branco só nas bordas INTERNAS (entre células), não em volta
        ix0 = x0 + (m if x0 > 0 else 0)
        iy0 = y0 + (m if y0 > 0 else 0)
        ix1 = x1 - (m if x1 < largura_total else 0)
        iy1 = y1 - (m if y1 < altura_ref else 0)
        cell_w, cell_h = max(1, ix1 - ix0), max(1, iy1 - iy0)
        cell_ratio = cell_w / cell_h
        im = Image.open(path).convert("RGB")
        w, h = im.size
        if (w / h) > cell_ratio:                 # foto mais larga que a célula: corta na horizontal
            new_w = int(h * cell_ratio)
            cx = _offset_corte_horizontal(path, w, new_w, foco=focos[i])
            im = im.crop((cx, 0, cx + new_w, h))
        else:                                    # mais alta: corta na vertical (folga acima das cabeças)
            new_h = int(w / cell_ratio)
            cy = _offset_corte_vertical(path, h, new_h, foco=focos[i])
            im = im.crop((0, cy, w, cy + new_h))
        im = im.resize((cell_w, cell_h))
        canvas.paste(im, (ix0, iy0))
    return canvas


LIMIAR_CASAMENTO_FOTO = 0.80  # abaixo disso, considera "não encontrado"


def match_photo_file(foto_arquivo, fotos_dir):
    """Casa o nome de foto do briefing com o arquivo real na pasta, tolerando
    diferenças de acento/espaço/maiúscula (via _normalize). Usa similaridade
    real (difflib) em vez de só prefixo em comum: nomes com o mesmo prefixo e
    o mesmo sufixo mas miolo diferente — típico dos cartazes, tipo
    'AGA_CARTAZ JM_Show OOP Nirvana_220626-KN' vs '...Orgulho LGBTQIA+...' —
    NÃO podem casar. Quando nada passa do limiar, retorna None (o chamador
    avisa 'não encontrado', em vez de trazer o arquivo errado calado)."""
    from difflib import SequenceMatcher
    if not foto_arquivo:
        return None
    alvo = _normalize(foto_arquivo)
    if not alvo:
        return None
    melhor, melhor_score = None, 0.0
    for cand in os.listdir(fotos_dir):
        cnorm = _normalize(cand)
        if not cnorm:
            continue
        score = SequenceMatcher(None, alvo, cnorm).ratio()
        # um nome curto contido inteiro no outro é um casamento forte
        if len(min(alvo, cnorm, key=len)) >= 6 and (alvo in cnorm or cnorm in alvo):
            score = max(score, 0.9)
        if score > melhor_score:
            melhor, melhor_score = cand, score
    if melhor and melhor_score >= LIMIAR_CASAMENTO_FOTO:
        return os.path.join(fotos_dir, melhor)
    return None


def resolve_foto(foto_arquivo, fotos_dir, target_ratio, warnings, contexto="",
                 foco=None, focos_unir=None, mosaico=False):
    """Resolve o campo 'foto_arquivo' (string única ou lista pra unir) num
    caminho de arquivo pronto pra usar em _replace_picture. `focos_unir` traz o
    foco POR foto da união (alinhado a foto_arquivo); onde a parte não tem foco
    próprio, cai no `foco` compartilhado da matéria (compat). `mosaico=True`
    compõe a união em mosaico (células variadas) em vez de fatias iguais."""
    if isinstance(foto_arquivo, list):
        paths = []
        for nome in foto_arquivo:
            p = match_photo_file(nome, fotos_dir)
            if not p:
                warnings.append(f"{contexto}: foto '{nome}' (parte de uma união) não encontrada.")
                return None
            paths.append(p)
        # foco por foto: usa o da parte; se a parte não tem, cai no foco da matéria
        lista_focos = [(focos_unir[i] if focos_unir and i < len(focos_unir) else None) or foco
                       for i in range(len(paths))]
        if mosaico:
            composta = _compose_mosaic_photo(paths, target_ratio, focos=lista_focos)
        else:
            composta = _compose_merged_photo(paths, target_ratio, focos=lista_focos)
        # nome determinístico (não depende do hash aleatório do Python), então
        # re-rodar sobrescreve em vez de acumular; a faxina em build_deck apaga
        # de vez depois do .pptx salvo. O foco entra na chave pra dois cortes
        # diferentes da mesma união não colidirem no mesmo arquivo.
        import hashlib
        assinatura = ("mosaico::" if mosaico else "") + "||".join(foto_arquivo) + "##" + "|".join(str(f) for f in lista_focos)
        chave = hashlib.md5(assinatura.encode("utf-8")).hexdigest()[:12]
        tmp_path = os.path.join(fotos_dir, f"_merged_{chave}.jpg")
        composta.save(tmp_path, quality=95)
        return tmp_path
    else:
        p = match_photo_file(foto_arquivo, fotos_dir)
        if not p:
            warnings.append(f"{contexto}: foto '{foto_arquivo}' não encontrada na pasta de fotos.")
        return p


# ---------------------------------------------------------------------------
# Populamento de slides
# ---------------------------------------------------------------------------

# Respiro (cm) entre a tag de editoria e a foto na capa, quando a editoria é
# comprida e a tag encosta na foto. Ver populate_capa.
FOLGA_TAG_FOTO_CM = 0.6

# Respiro (cm) entre a coluna de texto (corpo) e a foto na matéria de página
# única. O corpo NUNCA pode alargar por cima da foto — o teto de largura é a
# borda esquerda da foto menos essa folga. Ver populate_single_slide.
FOLGA_CORPO_FOTO_CM = 0.6


def populate_capa(slide, materia, fotos_dir, media_dir, warnings):
    textboxes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    pictures = [s for s in slide.shapes if s.shape_type == 13]

    title_box, body_box = None, None
    for tb in textboxes:
        runs = tb.text_frame.paragraphs[0].runs
        if not runs:
            continue
        if runs[0].font.bold and runs[0].font.size and runs[0].font.size >= 400000:
            title_box = tb
        elif not runs[0].font.bold:
            body_box = tb

    tag_box = None
    photo_box = None
    for p in pictures:
        try:
            fmt = p.image.ext
        except Exception:
            continue
        h_cm = p.height / 360000
        if fmt == "png" and 1.5 < h_cm < 2.2:
            tag_box = p
        elif fmt in ("jpg", "jpeg") and h_cm > 15:
            if photo_box is None or p.width * p.height > photo_box.width * photo_box.height:
                photo_box = p

    if title_box:
        pt_orig = _fonte_pt_do_titulo(title_box)
        # se dá pra tirar uma linha reduzindo a fonte só um pouquinho (o título
        # aproveita o espaço e fica com linhas mais cheias), reduz o necessário
        pt_titulo = _reduzir_fonte_titulo_menos_linhas(materia["titulo"], title_box.width, pt_orig)
        titulo_final = _evitar_viuva(materia["titulo"], title_box.width, pt_titulo)
        # rebalanceia as quebras (evita palavra órfã sozinha numa linha, tipo o
        # "CMG" da capa de Queiroz); mantém o mesmo nº de linhas
        titulo_final = _quebrar_titulo_balanceado(titulo_final, title_box.width, pt_titulo)
        _set_run_text_keep_format(title_box.text_frame, titulo_final)
        if pt_titulo != pt_orig:   # aplica a fonte reduzida (o corpo se reposiciona depois)
            for r in title_box.text_frame.paragraphs[0].runs:
                r.font.size = Emu(int(round(pt_titulo * 12700)))
    else:
        warnings.append("CAPA: não encontrei a caixa de título pra substituir.")
    if body_box:
        _set_run_text_keep_format(body_box.text_frame, _evitar_viuva_corpo(body_box, materia["corpo"]))
    else:
        warnings.append("CAPA: não encontrei a caixa de corpo pra substituir.")
    if title_box and body_box:
        limite_inferior = slide.part.package.presentation_part.presentation.slide_height - Emu(int(1.0 * 360000))
        # balanceador de massa: mantém o corpo em 19pt, distribui bem (estreitando
        # um tico a caixa — na capa o layout é fixo, não alarga) e compõe a última
        # linha no modelo PowerPoint (capa=True) evitando o toco tipo "(NR) nº22".
        _encaixar_corpo_preservando(title_box, body_box, materia["titulo"], materia["corpo"],
                                    limite_inferior, largura_max_corpo_emu=None,
                                    warnings=warnings, contexto="CAPA", capa=True)
    if tag_box:
        # a tag não pode invadir a foto (que na capa fica à direita): passa a
        # borda esquerda da foto, menos uma folga de respiro, como limite pra
        # tag encolher. 0.3cm deixava a tag "quase encostando" na foto quando a
        # editoria é comprida (ex.: "INOVAÇÃO E TECNOLOGIA"); 0.6cm dá um respiro
        # visível sem encolher a tag além do necessário.
        limite_tag = None
        if photo_box is not None and photo_box.left > tag_box.left:
            limite_tag = photo_box.left - Emu(int(FOLGA_TAG_FOTO_CM * 360000))
        _replace_tag(tag_box, materia["editoria"], limite_direito_emu=limite_tag)
    else:
        warnings.append("CAPA: não encontrei a tag de editoria pra substituir.")

    if photo_box:
        ratio = photo_box.width / photo_box.height
        foto_path = resolve_foto(materia.get("foto_arquivo"), fotos_dir, ratio, warnings, "CAPA",
                                 foco=materia.get("foco"), focos_unir=materia.get("focos_unir"),
                                 mosaico=materia.get("mosaico"))
        if foto_path:
            _replace_picture(photo_box, foto_path,
                             foto_inteira=materia.get("foto_inteira", False),
                             foco=materia.get("foco"), warnings=warnings, contexto="CAPA")
    else:
        warnings.append("CAPA: não encontrei a foto principal pra substituir.")


def populate_double_slide(slide, story_top, story_bottom, fotos_dir, media_dir, warnings):
    pictures = [s for s in slide.shapes if s.shape_type == 13]
    textboxes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]

    pictures.sort(key=_top_cm)
    textboxes.sort(key=_top_cm)

    photo_top, tag_top = sorted(pictures[:2], key=lambda s: s.width, reverse=True)
    photo_bottom, tag_bottom = sorted(pictures[2:], key=lambda s: s.width, reverse=True)

    title_top = body_top = title_bottom = body_bottom = None
    for tb in textboxes:
        runs = tb.text_frame.paragraphs[0].runs
        if not runs:
            continue
        bold = runs[0].font.bold
        if _top_cm(tb) < 16:
            if bold:
                title_top = tb
            else:
                body_top = tb
        else:
            if bold:
                title_bottom = tb
            else:
                body_bottom = tb

    slide_height = slide.part.package.presentation_part.presentation.slide_height
    limite_top_story = photo_top.top + photo_top.height  # corpo não passa da base da própria foto
    limite_bottom_story = min(photo_bottom.top + photo_bottom.height,
                               slide_height - Emu(int(1.0 * 360000)))

    for story, title_box, body_box, photo_box, tag_box, limite_inferior in [
        (story_top, title_top, body_top, photo_top, tag_top, limite_top_story),
        (story_bottom, title_bottom, body_bottom, photo_bottom, tag_bottom, limite_bottom_story),
    ]:
        qr_bottom = None
        if story.get("qr_path"):
            lado_qr = Emu(int(6.09 * 360000))
            margem_qr = Emu(int(2.0 * 360000))
            qr_bottom = title_box.top + lado_qr  # QR fica alinhado ao topo da tag/título
            limite_direito = slide.part.package.presentation_part.presentation.slide_width - lado_qr - margem_qr - Emu(int(0.3 * 360000))
            nova_largura = limite_direito - title_box.left
            if nova_largura < title_box.width:
                title_box.width = nova_largura
                # o corpo NÃO encolhe por causa do QR — ele fica só na linha
                # do título/tag, e o corpo (texto corrido, alinhado à esquerda)
                # não chega a alcançar essa área mesmo usando a largura cheia

        pt_titulo = _fonte_pt_do_titulo(title_box)
        # Matéria SEM QR: o corpo e a foto da dupla ficam à ESQUERDA, sobrando
        # espaço vazio à direita da caixa do título. Dar à caixa toda essa largura
        # livre faz o título respirar: um título comprido que, na largura do molde,
        # quebraria numa 3ª linha — caindo sobre o corpo, cujo topo é fixo — cabe
        # nas 2 linhas que o layout reserva (foi o ajuste manual do David no
        # 02-06 s3). Seguro por construção: mais largura NUNCA aumenta o nº de
        # linhas, o título é alinhado à esquerda (não se desloca) e a quebra segue
        # sendo a natural do PowerPoint. Com QR a caixa é estreitada logo abaixo
        # pra caber ao lado dele — por isso não se alarga aqui.
        if not story.get("qr_path"):
            slide_w = slide.part.package.presentation_part.presentation.slide_width
            max_titulo_w = slide_w - title_box.left - Emu(int(1.0 * 360000))
            if max_titulo_w > title_box.width:
                title_box.width = max_titulo_w
        titulo_final = _evitar_viuva(story["titulo"], title_box.width, pt_titulo)
        # reequilibra as quebras (evita linha curta no meio do título, tipo o
        # "voto é pela"), mantendo o mesmo nº de linhas e SEM mexer na largura
        titulo_final = _quebrar_titulo_balanceado(titulo_final, title_box.width, pt_titulo)
        _set_run_text_keep_format(title_box.text_frame, titulo_final)
        _set_run_text_keep_format(body_box.text_frame, _evitar_viuva_corpo(body_box, story["corpo"]))
        if story.get("qr_path"):
            _ajustar_titulo_largura_reduzida(title_box, titulo_final)
        # balanceador de massa: corpo em 19pt, distribuição sem vazão; a caixa pode
        # estreitar OU alargar (até a margem direita do slide, onde sobra espaço
        # abaixo do QR/título), e o título cede primeiro se faltar altura.
        slide_w = slide.part.package.presentation_part.presentation.slide_width
        largura_max = slide_w - body_box.left - Emu(int(1.0 * 360000))
        _encaixar_corpo_preservando(title_box, body_box, story["titulo"], story["corpo"],
                                    limite_inferior, largura_max_corpo_emu=largura_max,
                                    warnings=warnings, contexto=f"Matéria '{story['titulo']}'")
        _replace_tag(tag_box, story["editoria"])
        foto_path = resolve_foto(story.get("foto_arquivo"), fotos_dir,
                                  photo_box.width / photo_box.height, warnings,
                                  f"Matéria '{story['titulo']}'", foco=story.get("foco"),
                                  focos_unir=story.get("focos_unir"), mosaico=story.get("mosaico"))
        if foto_path:
            _replace_picture(photo_box, foto_path,
                             foto_inteira=story.get("foto_inteira", False),
                             foco=story.get("foco"), warnings=warnings,
                             contexto=f"Matéria '{story['titulo']}'")
        if story.get("qr_path"):
            _place_qr_code(slide, story["qr_path"], tag_box, slide.part.package.presentation_part.presentation.slide_width)


def populate_single_slide(slide, materia, fotos_dir, media_dir, warnings):
    """Matéria de página única (não-capa). Baseado no slide-modelo enviado
    pelo cliente (CEA/Centro de Memória): tag + título + corpo + foto.
    Caixas extras de informação (ex.: horários) não são geradas
    automaticamente — ficam como um ajuste manual pontual, se existirem."""
    pictures = [s for s in slide.shapes if s.shape_type == 13]
    textboxes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    textboxes.sort(key=_top_cm)

    on_slide_pics = [p for p in pictures if p.left >= 0]
    on_slide_pics_sorted = sorted(on_slide_pics, key=lambda p: p.width * p.height, reverse=True)
    photo_box = on_slide_pics_sorted[0] if on_slide_pics_sorted else None
    tag_box = None
    for p in on_slide_pics_sorted[1:]:
        h_cm = p.height / 360000
        if 1.5 < h_cm < 2.2:
            tag_box = p
            break

    title_box = textboxes[0] if textboxes else None
    body_box = textboxes[1] if len(textboxes) > 1 else None
    GAP_PADRAO = Emu(int(0.5 * 360000))

    if tag_box:
        novo_tag = _replace_tag(tag_box, materia["editoria"])
        # a tag sempre alinha com o topo da foto — hoje o motor sempre
        # remove as caixas extras de horário/info (não há dado do briefing
        # pra preenchê-las), então o slide final é sempre o caso "limpo".
        if photo_box:
            novo_tag.top = photo_box.top
            if title_box:
                title_box.top = novo_tag.top + novo_tag.height + GAP_PADRAO
    else:
        warnings.append(f"Página única '{materia['titulo']}': não encontrei a tag de editoria.")

    if title_box:
        pt_titulo = _fonte_pt_do_titulo(title_box)
        titulo_final = _evitar_viuva(materia["titulo"], title_box.width, pt_titulo)
        # reequilibra as quebras (linha curta no meio do título), mesmo nº de linhas
        titulo_final = _quebrar_titulo_balanceado(titulo_final, title_box.width, pt_titulo)
        _set_run_text_keep_format(title_box.text_frame, titulo_final)
    else:
        warnings.append(f"Página única '{materia['titulo']}': não encontrei a caixa de título.")
    if body_box:
        _set_run_text_keep_format(body_box.text_frame, _evitar_viuva_corpo(body_box, materia["corpo"]))
    else:
        warnings.append(f"Página única '{materia['titulo']}': não encontrei a caixa de corpo.")
    boxes_data = materia.get("boxes") or []
    linhas_slide = [s for s in slide.shapes if s.shape_type == 9]
    if title_box and body_box:
        limite_inferior = slide.part.package.presentation_part.presentation.slide_height - Emu(int(1.0 * 360000))
        # com boxes, o corpo (intro) não pode crescer até a área dos boxes:
        # o limite passa a ser o topo da 1ª linha divisória
        if boxes_data and linhas_slide:
            topo_boxes = min(s.top for s in linhas_slide if s.top is not None)
            limite_inferior = min(limite_inferior, topo_boxes - Emu(int(0.2 * 360000)))
        # TETO DE LARGURA do corpo (bug A + bug D). O corpo/intro pode alargar até
        # distribuir melhor a massa, MAS nunca por cima da foto: o teto é a borda
        # esquerda da foto menos FOLGA_CORPO_FOTO (col_max). Vale COM e SEM boxes:
        #  - sem boxes: impede o corpo de invadir a foto (16-12 s2, 07.07 s3,
        #    12-5 s2, 14-4 s2, 02-06 s2 — o molde de única "normal" traz a caixa de
        #    corpo em largura quase full-slide, ~37cm);
        #  - com boxes: DÁ espaço pro intro fechar em MENOS linhas e não colidir com
        #    o divisor (21-07 s5 "Confira" — antes ficava preso a 15,60 sem poder
        #    alargar e a 5ª linha caía sobre o traço do box).
        # Como o balanceador ancora o leque de larguras em largura0
        # (max(largura_max, largura0)), não basta passar um teto menor: encolhe-se a
        # caixa pro teto ANTES do balanceador. No-op onde a caixa já respeita a
        # coluna (a maioria dos moldes) → zero regressão.
        largura_max = None
        if photo_box is not None and photo_box.left is not None and photo_box.left > body_box.left:
            col_max = photo_box.left - body_box.left - Emu(int(FOLGA_CORPO_FOTO_CM * 360000))
            col_max = max(col_max, Emu(int(6.0 * 360000)))  # nunca estrangula demais
            largura_max = col_max
            if body_box.width > col_max:
                body_box.width = col_max
        elif not boxes_data:
            slide_w = slide.part.package.presentation_part.presentation.slide_width
            largura_max = slide_w - body_box.left - Emu(int(1.0 * 360000))
        _encaixar_corpo_preservando(title_box, body_box, materia["titulo"], materia["corpo"],
                                    limite_inferior, largura_max_corpo_emu=largura_max,
                                    warnings=warnings, contexto=f"Página única '{materia['titulo']}'")
    if photo_box:
        foto_path = resolve_foto(materia.get("foto_arquivo"), fotos_dir,
                                  photo_box.width / photo_box.height, warnings,
                                  f"Página única '{materia['titulo']}'", foco=materia.get("foco"),
                                  focos_unir=materia.get("focos_unir"), mosaico=materia.get("mosaico"))
        if foto_path:
            _replace_picture(photo_box, foto_path,
                             foto_inteira=materia.get("foto_inteira", False),
                             foco=materia.get("foco"), warnings=warnings,
                             contexto=f"Página única '{materia['titulo']}'")
    else:
        warnings.append(f"Página única '{materia['titulo']}': não encontrei a foto principal.")

    # QR code (opcional): o molde "única com QR" traz um placeholder quadrado na
    # posição definida no template; troca-se a imagem dele pelo QR real. A foto
    # e a tag são ignoradas na busca; a foto já é razão ~1.2 (não confunde com o
    # quadrado do QR). Sem QR pra essa matéria, um placeholder que sobrou é
    # removido.
    qr_path = materia.get("qr_path")
    achou_qr = _swap_qr_placeholder(slide, qr_path, ignorar=(tag_box,))
    if qr_path and not achou_qr:
        warnings.append(f"Página única '{materia['titulo']}': tem QR, mas não achei o "
                        f"placeholder de QR no molde (confira o template).")

    # boxes do briefing ("ABRIR BOX ...:"): se a matéria tem boxes, PREENCHE os
    # placeholders do layout (mantendo as linhas divisórias); se não tem,
    # remove as caixas extras e as linhas (comportamento padrão do layout limpo)
    placeholders = textboxes[2:]
    if boxes_data:
        _render_boxes(slide, body_box, placeholders, linhas_slide, boxes_data,
                      warnings, contexto=f"Página única '{materia['titulo']}'")
    else:
        for extra_tb in placeholders:
            extra_tb._element.getparent().remove(extra_tb._element)
        for line_shape in linhas_slide:
            line_shape._element.getparent().remove(line_shape._element)

    # a logo de encerramento só deve aparecer na ÚLTIMA matéria do JM inteiro —
    # como esse template une o slide-modelo com a logo já embutida, ela é
    # removida aqui sempre, e adicionada de volta (uma única vez) pelo
    # orquestrador (build_deck) depois de saber qual slide é realmente o último.
    for p in list(slide.shapes):
        if p.shape_type == 13:
            h_cm = p.height / 360000
            w_cm = p.width / 360000
            if 2.5 < w_cm < 6 and 2 < h_cm < 4.5 and p.left > 0:
                p._element.getparent().remove(p._element)


def _remover_assets_fora_da_prancheta(slide):
    """Remove o que está estacionado fora da área visível do slide e não deve
    ir pro arquivo final: o estoque de tags não usadas que o template carrega,
    e enfeites soltos na 'mesa' ao lado da prancheta (ex.: um quadrado laranja
    parado num left negativo). Trata qualquer tipo de shape — não só imagem —
    porque esses enfeites costumam ser auto-shapes (retângulos), que a versão
    antiga deixava passar. Só remove o que está 100% fora; um elemento que
    sangra pra fora mas ainda aparece em parte é preservado."""
    pres = slide.part.package.presentation_part.presentation
    sw, sh = pres.slide_width, pres.slide_height
    for shape in list(slide.shapes):
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h):
            continue
        totalmente_fora = (l + w <= 0) or (t + h <= 0) or (l >= sw) or (t >= sh)
        # comportamento antigo (mantido): imagem de estoque parada em left/top
        # negativo é removida mesmo que sangre um pouco pra dentro
        imagem_parada_fora = (shape.shape_type == 13 and (l < 0 or t < 0))
        if totalmente_fora or imagem_parada_fora:
            shape._element.getparent().remove(shape._element)


LOGO_FECHAMENTO_PATH = _os.path.join(TAGS_DIR, "logo_aga_fechamento.png")


def _adicionar_logo_final(slide, warnings=None):
    from pptx.util import Pt
    left = Emu(int(36.2 * 360000))
    top = Emu(int(25.22 * 360000))
    width = Emu(int(4.65 * 360000))
    height = Emu(int(3.1 * 360000))
    margem = Emu(int(0.6 * 360000))
    slide.shapes.add_picture(LOGO_FECHAMENTO_PATH, left, top, width, height)

    # dá "respiro" pra logo — MAS só estreita uma caixa se o TEXTO dela de fato
    # alcança a logo (uma linha que cai na faixa vertical da logo e chega perto
    # dela na horizontal). Antes o motor estreitava toda caixa cujo RETÂNGULO
    # passava por baixo da logo, mesmo quando o texto (alinhado à esquerda) nem
    # chegava lá — foi o que espremia o corpo e abria o vão do "busca". Quando
    # estreita mesmo, só re-gruda a viúva na largura nova (sem encolher a fonte).
    buffer_v = Emu(int(0.5 * 360000))
    MARG_INT = 91440
    logo_top, logo_bot = top, top + height
    for shape in list(slide.shapes):
        if not (shape.has_text_frame and shape.text_frame.text.strip()):
            continue
        if not (shape.left < left):
            continue
        if not (shape.top < logo_bot and shape.top + shape.height > logo_top - buffer_v):
            continue
        runs = shape.text_frame.paragraphs[0].runs
        if not runs:
            continue
        negrito = bool(runs[0].font.bold)
        pt = (runs[0].font.size / 12700) if runs[0].font.size else (40 if negrito else 19)
        texto = shape.text_frame.text
        usable = shape.width - 2 * MARG_INT
        line_h = int(pt * 12700 * 1.25)
        # mede o alcance horizontal de cada linha que cai na FAIXA da logo
        alcanca = False
        y = shape.top
        for seg in texto.split("\n"):
            for f in _fills_corpo(seg, shape.width, pt, negrito=negrito):
                if (y + line_h) > (logo_top - buffer_v) and y < logo_bot:
                    if shape.left + MARG_INT + int(f * usable) > (left - margem):
                        alcanca = True
                        break
                y += line_h
            if alcanca:
                break
        if not alcanca:
            continue                       # texto não chega na logo -> não mexe
        nova_largura = (left - margem) - shape.left
        if nova_largura <= 0 or nova_largura >= shape.width:
            continue
        shape.width = nova_largura
        if not negrito:                    # corpo: re-gruda a viúva na largura nova
            _set_run_text_keep_format(
                shape.text_frame,
                _glue_viuva_corpo(texto.replace("\xa0", " "), shape.width, pt))
            # ESTREITAR PELA LOGO AUMENTA O Nº DE LINHAS -> pode transbordar a base
            # do slide (bug C: o balanceador já tinha escolhido uma largura que
            # cabia, mas ela passava sob a logo; aqui a caixa encolhe e o texto
            # cresce em altura). O corpo então cede a fonte (19->18->17pt, piso de
            # leitura) até caber na vertical; se nem no piso couber, avisa em vez de
            # deixar vazar mudo. Casos: 02-06 s4 "Dicas".
            # Mede o FUNDO real do texto na MESMA régua do renderizador (fator de
            # linha 1,2 + inset da caixa), pra só baixar a fonte quando de fato
            # ultrapassa a base do slide — senão derruba corpo que ainda cabia
            # (ex.: 30-06 s4 "Semear", que a 19pt fecha em 29,18 < 29,7).
            slide_h = slide.part.package.presentation_part.presentation.slide_height
            limite = slide_h - Emu(int(0.2 * 360000))
            def _fundo(p):
                n = _contar_linhas_texto(shape.text_frame.text,
                                         int(shape.width * FATOR_WRAP_RENDER), p, negrito=False)
                return shape.top + Emu(int(0.13 * 360000)) + n * int(p * 12700 * 1.2)
            pt_novo = pt
            for cand in (pt, 18.0, PT_PISO_CORPO):
                if cand > pt:
                    continue
                if _fundo(cand) <= limite:
                    pt_novo = cand
                    break
            else:
                pt_novo = PT_PISO_CORPO
            if abs(pt_novo - pt) > 0.01:
                for paragraph in shape.text_frame.paragraphs:
                    for r in paragraph.runs:
                        r.font.size = Pt(pt_novo)
                _set_run_text_keep_format(
                    shape.text_frame,
                    _glue_viuva_corpo(texto.replace("\xa0", " "), shape.width, pt_novo))
            if _fundo(pt_novo) > limite and warnings is not None:
                warnings.append(
                    f"Matéria de rodapé: muito texto pro espaço ao lado da logo — "
                    f"corpo mantido em {int(pt_novo)}pt e ainda pode transbordar; "
                    f"considere encurtar o texto ou ajustar à mão.")


# ---------------------------------------------------------------------------
# Orquestrador principal
# ---------------------------------------------------------------------------

_PADRAO_MERGED = re.compile(r"^_merged_[0-9a-f]+\.jpg$", re.IGNORECASE)


def _limpar_merged(fotos_dir):
    """Apaga os JPEGs intermediários de fotos unidas (`_merged_*.jpg`) que o
    motor grava na pasta da semana ao compor uma união. Depois do .pptx salvo,
    os bytes já ficam embutidos no arquivo — esses intermediários não servem
    mais. Rodado no início do build (limpa restos de execuções anteriores) e no
    fim (limpa os desta execução), pra não acumular na pasta do cliente."""
    try:
        for nome in os.listdir(fotos_dir):
            if _PADRAO_MERGED.match(nome):
                try:
                    os.remove(os.path.join(fotos_dir, nome))
                except OSError:
                    pass
    except OSError:
        pass


LIMIAR_AREA_FOTO_CM2 = 100.0  # área mínima (cm²) pra uma imagem contar como
                              # "foto grande" (dupla ~190, única ~386, cartaz
                              # ~1247); tag/QR/logo ficam bem abaixo.


def _classificar_slide_template(slide):
    """Descobre o PAPEL de um slide-modelo pela assinatura das formas (não pela
    posição no arquivo): 'cartaz' (1 foto full-bleed, sem texto), 'dupla' (2+
    fotos grandes), 'unica_qr' (1 foto grande + placeholder quadrado do QR),
    'unica' (1 foto grande), ou None (não reconhecido). Assim, inserir ou
    reordenar páginas no template não quebra mais o motor."""
    pics = [s for s in slide.shapes if s.shape_type == 13
            and s.left is not None and s.left >= 0]
    tem_texto = any(s.has_text_frame and s.text_frame.text.strip() for s in slide.shapes)
    if not pics:
        return None

    def area_cm(p):
        return (p.width / 360000.0) * (p.height / 360000.0)

    def ratio(p):
        return (p.width / p.height) if p.height else 0

    if len(pics) == 1 and not tem_texto:
        return "cartaz"
    grandes = [p for p in pics if area_cm(p) >= LIMIAR_AREA_FOTO_CM2]
    if len(grandes) >= 2:
        return "dupla"
    if len(grandes) == 1:
        for p in pics:
            if p is grandes[0]:
                continue
            lado_cm = p.width / 360000.0
            if 0.85 <= ratio(p) <= 1.18 and 3.0 <= lado_cm <= 8.0:
                return "unica_qr"   # tem o quadrado do QR
        return "unica"
    return None


def _detectar_papeis_template(slides):
    """Mapeia os slides-modelo do template pelo papel de cada um (ver
    _classificar_slide_template). Capa = 1º slide; página final = último; o
    miolo é classificado. Retorna os índices de cada molde — nada é chumbado
    por número, então editar o template (ex.: incluir 'única com QR') não
    desalinha os cartazes nem exige mexer no código."""
    n = len(slides)
    papeis = {"capa": 0, "final": n - 1, "dupla": None, "unica": None,
              "unica_qr": None, "cartaz": [], "materia_templates": []}
    for i in range(1, n - 1):   # exclui capa (0) e página final (último)
        papel = _classificar_slide_template(slides[i])
        if papel == "cartaz":
            papeis["cartaz"].append(i)
        elif papel in ("dupla", "unica", "unica_qr"):
            if papeis[papel] is None:
                papeis[papel] = i
            papeis["materia_templates"].append(i)
    papeis["cartaz"] = tuple(papeis["cartaz"])
    papeis["materia_templates"] = tuple(papeis["materia_templates"])
    return papeis


# ---------------------------------------------------------------------------
# Multi-unidade: tags de unidade (MG/QZ/LM/CDS/GO) — selos nos slides + um
# gabarito (página final) por conjunto de slides. Ver CEREBRO_DO_JM.md.
# ---------------------------------------------------------------------------

# Uma cor fixa por sigla, pra o selo do slide bater com o rodapé do gabarito.
# MG = laranja da marca; as demais escolhidas pra contraste entre si.
CORES_UNIDADE = {
    "MG": (0xF2, 0x65, 0x22),   # laranja da marca
    "QZ": (0x1F, 0x6F, 0xEB),   # azul
    "LM": (0x2E, 0x9E, 0x5B),   # verde
    "CDS": (0x7E, 0x3F, 0xF2),  # roxo
    "GO": (0x8E, 0x24, 0x4D),   # vinho
}


def _agrupar_gabaritos(content_units):
    """`content_units`: lista, na ORDEM dos slides de conteúdo, do conjunto de
    unidades de cada slide. Devolve os grupos de gabarito — unidades que
    recebem EXATAMENTE o mesmo conjunto de slides compartilham um gabarito.
    Cada grupo: {'unidades': [...], 'slide_indices': [...]}. Ordem preservada
    pela 1ª aparição da unidade (e dos grupos)."""
    ordem = []
    for us in content_units:
        for u in us:
            if u not in ordem:
                ordem.append(u)
    idx_por_unidade = {
        u: tuple(i for i, us in enumerate(content_units) if u in us)
        for u in ordem
    }
    grupos = []
    for u in ordem:
        chave = idx_por_unidade[u]
        for g in grupos:
            if g["_chave"] == chave:
                g["unidades"].append(u)
                break
        else:
            grupos.append({"_chave": chave, "unidades": [u],
                           "slide_indices": list(chave)})
    for g in grupos:
        g.pop("_chave")
    return grupos


def _estilizar_circulo_unidade(circ, sig, pt):
    """Preenche um oval com a cor da unidade e escreve a sigla em branco no
    centro (usado no selo do slide e no rodapé do gabarito)."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(*CORES_UNIDADE.get(sig, (0x5A, 0x5A, 0x5A)))
    circ.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    circ.line.width = Emu(int(0.03 * 360000))
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = sig
    run.font.size = Pt(pt if len(sig) <= 2 else pt - 2)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_selo_unidades(slide, unidades, prs):
    """Pilha vertical de círculos coloridos (um por unidade, com a sigla
    dentro) encostada na borda DIREITA, FORA da prancheta, alinhada ao topo.
    Conteúdo fora do slide não sai no render/PDF/impressão — o selo é só um
    guia visual pra quem monta o JM no PowerPoint."""
    from pptx.enum.shapes import MSO_SHAPE
    # círculo dimensionado pra sigla legível (fonte 18) — uma pessoa lê isso pra
    # saber pra onde vai a matéria, então precisa dar leitura.
    diam = Emu(int(1.35 * 360000))
    gap = Emu(int(0.28 * 360000))
    left = prs.slide_width + Emu(int(0.35 * 360000))
    top = Emu(int(0.50 * 360000))
    for sig in unidades:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diam, diam)
        _estilizar_circulo_unidade(circ, sig, 18)
        top = top + diam + gap


def _add_rodape_gabarito(slide, unidades, prs):
    """Fileira de círculos coloridos com as siglas, centralizada no rodapé do
    gabarito — a versão visual das tags que no briefing ficam no fim de cada
    gabarito. Fica DENTRO da prancheta (aparece no arquivo/impressão)."""
    from pptx.enum.shapes import MSO_SHAPE
    # tira a tag manual que o molde do template já traz no rodapé (uma caixa de
    # texto com só a sigla, ex.: 'MG') — vira o círculo colorido no lugar.
    for shape in list(slide.shapes):
        try:
            if (shape.has_text_frame
                    and shape.text_frame.text.strip().upper() in CORES_UNIDADE):
                shape._element.getparent().remove(shape._element)
        except Exception:
            pass
    # mesmo tamanho do selo do slide (fonte 18 / ~1.35cm), pra bater visualmente.
    diam = Emu(int(1.35 * 360000))
    gap = Emu(int(0.35 * 360000))
    n = len(unidades)
    total_w = diam * n + gap * (n - 1)
    left = (prs.slide_width - total_w) // 2
    top = prs.slide_height - Emu(int(1.75 * 360000))
    for sig in unidades:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diam, diam)
        _estilizar_circulo_unidade(circ, sig, 18)
        left = left + diam + gap


def build_deck(template_path, paginas, fotos_dir, output_path, media_dir=None,
               cartazes=None):
    warnings = []
    _limpar_merged(fotos_dir)   # tira restos de fotos unidas de execuções anteriores
    if media_dir is None:
        media_dir = extract_media_from_template(template_path)
    prs = Presentation(template_path)
    slides_originais = list(prs.slides)

    # descobre qual slide é qual pelo conteúdo (não pela posição) — ver
    # _detectar_papeis_template. Índice nenhum fica chumbado.
    papeis = _detectar_papeis_template(slides_originais)
    capa_idx = papeis["capa"]
    if papeis["dupla"] is None:
        warnings.append("Template: não achei o molde de matéria DUPLA (2 fotos).")
    if papeis["unica"] is None:
        warnings.append("Template: não achei o molde de matéria ÚNICA (1 foto).")

    # unidades por página. multi_unidade liga SÓ quando alguma página trouxe
    # tag explícita no briefing; sem tag, tudo é MG e o comportamento é o de
    # sempre (1 gabarito, sem selos) — zero regressão nas edições antigas.
    multi_unidade = any(p.get("unidades") for p in paginas)
    units_por_pagina = [(p.get("unidades") or ["MG"]) for p in paginas]
    todas_unidades = []
    for us in units_por_pagina:
        for u in us:
            if u not in todas_unidades:
                todas_unidades.append(u)

    capa_materia = paginas[0]["materias"][0]
    populate_capa(slides_originais[capa_idx], capa_materia, fotos_dir, media_dir, warnings)
    _remover_assets_fora_da_prancheta(slides_originais[capa_idx])

    template_dupla = slides_originais[papeis["dupla"]] if papeis["dupla"] is not None else None
    template_unica = slides_originais[papeis["unica"]] if papeis["unica"] is not None else None
    template_unica_qr = slides_originais[papeis["unica_qr"]] if papeis["unica_qr"] is not None else None
    novas_materia_slides = []

    for pagina in paginas[1:]:
        if pagina["tipo"] == "dupla":
            m1, m2 = (pagina["materias"] + [None])[:2]
            new_slide = duplicate_slide(prs, template_dupla)
            populate_double_slide(new_slide, m1, m2 or {"editoria": m1["editoria"], "titulo": "", "corpo": "", "foto_arquivo": None},
                                   fotos_dir, media_dir, warnings)
            _remover_assets_fora_da_prancheta(new_slide)
            novas_materia_slides.append(new_slide)
        else:
            materia = pagina["materias"][0]
            # página única COM QR usa o molde que tem o placeholder do QR (se o
            # template tiver esse molde); sem QR, o molde simples. Cai no
            # simples se não houver molde de QR (degradação graciosa).
            usa_molde_qr = materia.get("qr_path") and template_unica_qr is not None
            molde = template_unica_qr if usa_molde_qr else template_unica
            new_slide = duplicate_slide(prs, molde)
            populate_single_slide(new_slide, materia, fotos_dir, media_dir, warnings)
            _remover_assets_fora_da_prancheta(new_slide)
            novas_materia_slides.append(new_slide)

    novos_cartaz_slides = []
    # logo de fechamento: vai nas páginas marcadas com "/ FECHA COM LOGO DA AGA /"
    # no briefing — no multi-unidade há uma última página por unidade, cada uma
    # com a sua logo. Sem marcador nenhum (fallback), mantém o comportamento de
    # sempre: logo na última matéria (ou na capa se não houver matéria).
    paginas_com_logo = [i for i, p in enumerate(paginas) if p.get("fecha_com_logo")]
    if paginas_com_logo:
        for i in paginas_com_logo:
            slide = slides_originais[capa_idx] if i == 0 else novas_materia_slides[i - 1]
            _adicionar_logo_final(slide, warnings)
    elif novas_materia_slides:
        _adicionar_logo_final(novas_materia_slides[-1], warnings)
    else:
        _adicionar_logo_final(slides_originais[capa_idx], warnings)

    if cartazes and papeis["cartaz"]:
        template_cartaz = slides_originais[papeis["cartaz"][0]]
        for img_path in cartazes:
            new_slide = duplicate_slide(prs, template_cartaz)
            pic = [s for s in new_slide.shapes if s.shape_type == 13][0]
            left, top, w, h = pic.left, pic.top, pic.width, pic.height
            old_el = pic._element
            parent = old_el.getparent()
            pos = list(parent).index(old_el)
            parent.remove(old_el)
            new_pic = new_slide.shapes.add_picture(img_path, left, top, w, h)
            new_el = new_pic._element
            new_el.getparent().remove(new_el)
            parent.insert(pos, new_el)
            novos_cartaz_slides.append(new_slide)

    # --- gabaritos: agrupa unidades por conjunto de slides ---
    # ordem dos slides de CONTEÚDO no deck final: capa, matérias, cartazes.
    # Cada matéria-slide herda as unidades da sua página; cartaz vai pra TODAS.
    final_idx = papeis["final"]
    if multi_unidade:
        content_units = ([units_por_pagina[0]]
                         + units_por_pagina[1:]
                         + [list(todas_unidades)] * len(novos_cartaz_slides))
        grupos = _agrupar_gabaritos(content_units)
    else:
        content_units = None
        grupos = None
    n_gab = len(grupos) if grupos else 1
    # o template traz 1 molde de página final; pra N gabaritos, duplica N-1.
    gabarito_dups = [duplicate_slide(prs, slides_originais[final_idx])
                     for _ in range(n_gab - 1)]

    # --- reordenar: capa, novas matérias, novos cartazes, gabaritos; remove antigos ---
    sld_id_lst = prs.slides._sldIdLst
    all_ids = list(sld_id_lst)

    capa_id = all_ids[capa_idx]
    final_molde_id = all_ids[final_idx]
    # remove TODO o miolo do template (tudo que não é a capa nem a página
    # final): esses slides são só moldes — os slides novos gerados entram no
    # lugar. Só entre os slides ORIGINAIS (os novos ficam no fim de all_ids e
    # não podem entrar na remoção). Definir por exclusão (capa + final) é
    # robusto a qualquer edição do template, classificada ou não.
    n_orig = len(slides_originais)
    indices_a_remover = {i for i in range(n_orig) if i not in (capa_idx, final_idx)}
    old_ids_a_remover = [all_ids[i] for i in indices_a_remover]

    # os novos ids ficam no FIM de all_ids, na ordem de criação: matérias,
    # cartazes e (por último) as duplicatas de gabarito.
    n_mat = len(novas_materia_slides)
    n_cart = len(novos_cartaz_slides)
    n_dup = len(gabarito_dups)
    n_novos = n_mat + n_cart + n_dup
    tail = all_ids[len(all_ids) - n_novos:] if n_novos else []
    novos_materia_ids = tail[:n_mat]
    novos_cartaz_ids = tail[n_mat:n_mat + n_cart]
    dup_ids = tail[n_mat + n_cart:]
    gabarito_ids = [final_molde_id] + dup_ids   # N no total, na ordem de `grupos`

    for sid in all_ids:
        sld_id_lst.remove(sid)

    ordem_final = [capa_id] + novos_materia_ids + novos_cartaz_ids + gabarito_ids
    for sid in ordem_final:
        sld_id_lst.append(sid)

    for sid in old_ids_a_remover:
        # remove definitivamente os slides antigos (matérias e cartazes) não usados nesta edição
        rId = sid.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        prs.part.drop_rel(rId)

    # --- selos (fora do slide) nos slides de conteúdo + rodapé nos gabaritos ---
    if multi_unidade:
        content_slides = ([slides_originais[capa_idx]]
                          + novas_materia_slides + novos_cartaz_slides)
        for slide, us in zip(content_slides, content_units):
            _add_selo_unidades(slide, us, prs)
        gabarito_slides = [slides_originais[final_idx]] + gabarito_dups
        for slide, grupo in zip(gabarito_slides, grupos):
            _add_rodape_gabarito(slide, grupo["unidades"], prs)

    prs.save(output_path)
    _limpar_merged(fotos_dir)   # os intermediários já estão embutidos no .pptx
    return warnings, grupos


def _encontrar_soffice():
    """Procura o LibreOffice tanto no PATH quanto nos lugares onde ele é
    instalado no Mac e no Windows, porque o app quase nunca coloca o
    comando `soffice` no PATH sozinho — quem instala pelo instalador
    normal (ou pelo Homebrew cask) fica com o executável só dentro do
    pacote do aplicativo."""
    import shutil
    cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if cmd:
        return cmd
    candidatos = [
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        _os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
        "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
    ]
    for caminho in candidatos:
        if _os.path.isfile(caminho):
            return caminho
    return None


def _pdf_para_imagens(pdf_path, tmp_dir, dpi=100):
    """Converte cada página do PDF numa imagem, sem depender de nenhum
    programa externo: usa o PyMuPDF (pacote `pymupdf`, instalado via pip e
    que já traz o motor de renderização embutido). Se por algum motivo o
    PyMuPDF não estiver disponível, cai pro Poppler (`pdftoppm`) caso ele
    exista na máquina."""
    import subprocess
    import shutil

    try:
        import fitz  # PyMuPDF
        saidas = []
        with fitz.open(pdf_path) as doc:
            zoom = dpi / 72.0
            matriz = fitz.Matrix(zoom, zoom)
            for i, pagina in enumerate(doc):
                pix = pagina.get_pixmap(matrix=matriz)
                destino = _os.path.join(tmp_dir, f"pg-{i + 1:03d}.png")
                pix.save(destino)
                saidas.append(destino)
        return saidas

    except ImportError:
        pdftoppm_cmd = shutil.which("pdftoppm")
        if not pdftoppm_cmd:
            raise RuntimeError(
                "Falta uma forma de transformar o PDF em imagens. Instale o "
                "PyMuPDF com: pip3 install pymupdf"
            )
        subprocess.run(
            [pdftoppm_cmd, "-png", "-r", str(dpi), pdf_path,
             _os.path.join(tmp_dir, "pg")],
            check=True, capture_output=True,
        )
        return sorted(
            _os.path.join(tmp_dir, f) for f in _os.listdir(tmp_dir)
            if f.startswith("pg")
        )


def _montar_grade_miniaturas(final_slide, thumb_paths, slide_w, slide_h,
                             reservar_rodape=False):
    """Limpa as imagens antigas do slide-gabarito e monta a grade de
    miniaturas centralizada. `reservar_rodape=True` deixa uma faixa livre
    embaixo pros círculos de unidade do rodapé (gabarito multi-unidade)."""
    for shape in list(final_slide.shapes):
        if shape.shape_type == 13 or shape.shape_type == 6:  # picture ou group
            shape._element.getparent().remove(shape._element)

    margem = Emu(int(0.5 * 360000))
    rodape = Emu(int(2.2 * 360000)) if reservar_rodape else Emu(0)
    area_w = slide_w - 2 * margem
    area_h = slide_h - 2 * margem - rodape

    n = len(thumb_paths)
    if n == 0:
        return
    cols = math.ceil(math.sqrt(n * (area_w / area_h)))
    cols = max(1, min(cols, n))
    rows = math.ceil(n / cols)

    gap = Emu(int(0.3 * 360000))
    cell_w = (area_w - gap * (cols - 1)) // cols
    thumb_ratio = slide_w / slide_h
    cell_h = int(cell_w / thumb_ratio)
    if rows * cell_h + gap * (rows - 1) > area_h:
        cell_h = (area_h - gap * (rows - 1)) // rows
        cell_w = int(cell_h * thumb_ratio)

    grid_w = cell_w * cols + gap * (cols - 1)
    grid_h = cell_h * rows + gap * (rows - 1)
    start_x = (slide_w - grid_w) // 2
    start_y = margem + (area_h - grid_h) // 2

    for i, thumb in enumerate(thumb_paths):
        r, c = divmod(i, cols)
        left = start_x + c * (cell_w + gap)
        top = start_y + r * (cell_h + gap)
        final_slide.shapes.add_picture(thumb, left, top, cell_w, cell_h)


def regenerate_final_page(pptx_path, gabaritos=None, soffice_script_path=None):
    """Renderiza todos os slides já gerados (exceto os próprios gabaritos) e
    monta uma grade de miniaturas em cada gabarito (página final), substituindo
    o que estava lá antes (que pertencia a outra edição).

    `gabaritos` (multi-unidade): lista de {'unidades', 'slide_indices'} — um
    gabarito por grupo, cada um com as miniaturas SÓ dos seus slides. None
    (padrão) = um único gabarito com as miniaturas de todos os slides.

    Precisa do LibreOffice instalado (para renderizar o .pptx tal como ele
    aparece). A conversão de PDF pra imagem é feita pelo PyMuPDF, que vem
    junto via pip — não precisa de Poppler. Se o LibreOffice não for
    encontrado, a função avisa e não faz nada: o resto do arquivo (capa,
    matérias, cartazes) já foi gerado normalmente antes desse passo."""
    import subprocess
    import tempfile

    soffice_cmd = _encontrar_soffice()
    if not soffice_cmd:
        return ("Página final (miniaturas) não gerada: LibreOffice não "
                "encontrado — instale-o e rode de novo. O restante do JM "
                "(capa, matérias, cartazes) saiu normalmente.")

    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    n_gab = len(gabaritos) if gabaritos else 1
    n_content = len(slides) - n_gab   # slides de conteúdo (antes dos gabaritos)

    tmp_dir = tempfile.mkdtemp(prefix="jm_final_")
    # perfil próprio evita conflito com uma instância do LibreOffice já
    # aberta pelo usuário (aí o --headless não trava esperando a outra).
    # A URL do perfil TEM que ser montada com Path.as_uri(): no Windows o
    # caminho é C:\Users\... e um "file://" + caminho cru vira uma URL
    # malformada (file://C:\Users\...) que faz o LibreOffice engasgar no
    # start (erro "bootstrap.ini corrompido"). O as_uri() gera file:///C:/...
    # no Windows e file:///Users/... no Mac — correto nos dois.
    from pathlib import Path
    perfil_url = Path(_os.path.join(tmp_dir, "lo_profile")).as_uri()
    try:
        resultado = subprocess.run(
            [soffice_cmd, "--headless", "--norestore", "--nofirststartwizard",
             f"-env:UserInstallation={perfil_url}",
             "--convert-to", "pdf", "--outdir", tmp_dir, pptx_path],
            capture_output=True, timeout=180,
        )
    except subprocess.TimeoutExpired:
        # LibreOffice travou (ex.: instalação com problema abrindo um diálogo).
        # Não deixa o JM inteiro pendurado: avisa e segue sem a página final.
        return ("Página final (miniaturas) não gerada: o LibreOffice demorou "
                "demais (a instalação pode estar com problema — tente reparar/"
                "reinstalar). O restante do JM saiu normalmente.")
    generated_pdf = os.path.join(tmp_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")
    if resultado.returncode != 0 or not os.path.isfile(generated_pdf):
        return ("Página final (miniaturas) não gerada: o LibreOffice não "
                "conseguiu converter o arquivo (a instalação pode estar com "
                "problema — tente reparar/reinstalar). O restante do JM saiu "
                "normalmente.")
    # uma imagem por slide; os slides de conteúdo são os n_content primeiros
    # (os gabaritos ficam no fim e não viram miniatura de si mesmos).
    thumbs_todos = _pdf_para_imagens(generated_pdf, tmp_dir, dpi=100)
    content_thumbs = thumbs_todos[:n_content]

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if gabaritos:
        # um gabarito por grupo: cada um recebe só as miniaturas dos seus slides.
        for k, grupo in enumerate(gabaritos):
            gab_slide = slides[n_content + k]
            sel = [content_thumbs[i] for i in grupo["slide_indices"]
                   if i < len(content_thumbs)]
            _montar_grade_miniaturas(gab_slide, sel, slide_w, slide_h,
                                     reservar_rodape=True)
    else:
        # comportamento clássico: um único gabarito com TODAS as miniaturas.
        _montar_grade_miniaturas(slides[-1], content_thumbs, slide_w, slide_h)

    prs.save(pptx_path)


def limpar_coautoria(pptx_path):
    """Remove do arquivo final as partes de rastreamento de coautoria/
    alterações do PowerPoint (authors.xml, revisionInfo.xml, changesInfos/).

    Por que isso é necessário: essas partes guardam o histórico de quem
    mexeu em cada slide, identificando os slides por um ID interno. Como o
    motor remove os slides originais do template e cria os novos, esses IDs
    deixam de existir — e aí o PowerPoint abre o arquivo, vê referências
    apontando pra slides que sumiram e acusa o arquivo como corrompido
    ('o PowerPoint encontrou um problema... deseja reparar?'). Num arquivo
    gerado do zero essas partes não têm utilidade nenhuma, então a saída
    mais limpa é simplesmente removê-las. Retorna quantas partes tirou."""
    import posixpath

    def eh_coautoria(nome):
        return (nome in ("ppt/authors.xml", "ppt/revisionInfo.xml")
                or nome.startswith("ppt/changesInfos/"))

    with zipfile.ZipFile(pptx_path) as z:
        nomes = z.namelist()
        conteudos = {n: z.read(n) for n in nomes}

    dropar = set(n for n in nomes if eh_coautoria(n))
    # inclui também os _rels dessas partes, se existirem
    for n in nomes:
        if n.endswith(".rels"):
            base = n.replace("_rels/", "").rsplit(".rels", 1)[0]
            if eh_coautoria(base):
                dropar.add(n)
    if not dropar:
        return 0

    # tira as <Relationship> que apontam pras partes removidas
    rels_name = "ppt/_rels/presentation.xml.rels"
    if rels_name in conteudos:
        rels = conteudos[rels_name].decode("utf-8")

        def rel_dropada(m):
            alvo = re.search(r'Target="([^"]+)"', m.group(0))
            if not alvo:
                return False
            full = posixpath.normpath(posixpath.join("ppt", alvo.group(1)))
            return full in dropar

        rels = re.sub(r"<Relationship\b[^>]*/>",
                      lambda m: "" if rel_dropada(m) else m.group(0), rels)
        conteudos[rels_name] = rels.encode("utf-8")

    # tira os <Override> do Content_Types das partes removidas
    ct_name = "[Content_Types].xml"
    if ct_name in conteudos:
        ct = conteudos[ct_name].decode("utf-8")

        def override_dropado(m):
            pn = re.search(r'PartName="([^"]+)"', m.group(0))
            return bool(pn) and pn.group(1).lstrip("/") in dropar

        ct = re.sub(r"<Override\b[^>]*/>",
                    lambda m: "" if override_dropado(m) else m.group(0), ct)
        conteudos[ct_name] = ct.encode("utf-8")

    tmp = pptx_path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for n in nomes:
            if n not in dropar:
                out.writestr(n, conteudos[n])
    os.replace(tmp, pptx_path)
    return len(dropar)
